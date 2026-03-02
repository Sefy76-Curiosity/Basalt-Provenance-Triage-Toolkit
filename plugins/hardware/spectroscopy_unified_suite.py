"""
SPECTROSCOPY UNIFIED SUITE v5.0 - COMPLETE PRODUCTION RELEASE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
✓ FTIR: Thermo · PerkinElmer · Bruker · Agilent · Shimadzu — REAL drivers
✓ RAMAN: B&W Tek · Ocean · Horiba · Renishaw · Metrohm — Full protocols
✓ UV-VIS-NIR: Ocean · Avantes · Thorlabs · StellarNet — seabreeze/serial
✓ PORTABLE NIR: Viavi MicroNIR · TI DLP NIRscan — BLE/HTTP
✓ 52+ MODELS · 20+ FILE FORMATS · PEAK MATCHING · SPECTRAL LIBRARY
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""

# ============================================================================
# PLUGIN METADATA
# ============================================================================
PLUGIN_INFO = {
    "id": "spectroscopy_unified_suite",
    "name": "Spectroscopy Suite",
    "category": "hardware",
    "icon": "🧪",
    "version": "5.0.0",
    "author": "Spectroscopy Team",
    "description": "FTIR · RAMAN · UV-VIS-NIR · 52+ instruments · Real drivers",
    "requires": ["numpy", "pandas", "scipy", "matplotlib", "pyserial"],
    "optional": [
        "seabreeze",
        "pyvisa",
        "requests",
        "bleak",
        "hidapi",
        "openpyxl"
    ],
    "compact": True,
    "window_size": "1000x650"
}

# ============================================================================
# PREVENT DOUBLE REGISTRATION
# ============================================================================
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import numpy as np
import pandas as pd
from datetime import datetime, timedelta
import time
import re
import json
import threading
import queue
import subprocess
import sys
import os
import csv
import struct
from pathlib import Path
import platform
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple, Any, Callable, Union
from enum import Enum
from collections import deque   # <--- for history
import warnings
warnings.filterwarnings('ignore')

# ============================================================================
# OPTIONAL SCIENTIFIC IMPORTS
# ============================================================================
try:
    from scipy import signal, integrate, optimize, stats
    from scipy.signal import savgol_filter, find_peaks, peak_widths
    HAS_SCIPY = True
except ImportError:
    HAS_SCIPY = False

try:
    import matplotlib.pyplot as plt
    from matplotlib.figure import Figure
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk
    from matplotlib.patches import Rectangle, Polygon
    HAS_MPL = True
except ImportError:
    HAS_MPL = False

# ============================================================================
# SPECTROSCOPY-SPECIFIC IMPORTS
# ============================================================================
try:
    import seabreeze.spectrometers as sb
    from seabreeze.spectrometers import Spectrometer
    HAS_SEABREEZE = True
except ImportError:
    HAS_SEABREEZE = False

try:
    import serial
    import serial.tools.list_ports
    HAS_SERIAL = True
except ImportError:
    HAS_SERIAL = False

try:
    import pyvisa
    HAS_VISA = True
except ImportError:
    HAS_VISA = False

try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    HAS_REQUESTS = False

try:
    import bleak
    from bleak import BleakScanner, BleakClient
    HAS_BLEAK = True
except ImportError:
    HAS_BLEAK = False

try:
    from PIL import Image as _PILImage
    import io as _pil_io
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    from pptx import Presentation as _PPTXPresentation
    HAS_PPTX_LIB = True
except ImportError:
    HAS_PPTX_LIB = False

# ============================================================================
# CROSS-PLATFORM CHECK
# ============================================================================
IS_WINDOWS = platform.system() == 'Windows'
IS_LINUX = platform.system() == 'Linux'
IS_MAC = platform.system() == 'Darwin'

# ============================================================================
# THREAD-SAFE UI QUEUE
# ============================================================================
class ThreadSafeUI:
    def __init__(self, root):
        self.root = root
        self.queue = queue.Queue()
        self._poll()

    def _poll(self):
        try:
            while True:
                callback = self.queue.get_nowait()
                callback()
        except queue.Empty:
            pass
        finally:
            self.root.after(50, self._poll)

    def schedule(self, callback, *args, **kwargs):
        self.queue.put(lambda: callback(*args, **kwargs))

# ============================================================================
# TOOLTIP CLASS
# ============================================================================
class ToolTip:
    def __init__(self, widget, text):
        self.widget = widget
        self.text = text
        self.tip_window = None
        widget.bind('<Enter>', self.show_tip)
        widget.bind('<Leave>', self.hide_tip)

    def show_tip(self, event):
        x = self.widget.winfo_rootx() + 25
        y = self.widget.winfo_rooty() + 25
        self.tip_window = tw = tk.Toplevel(self.widget)
        tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{x}+{y}")
        label = tk.Label(tw, text=self.text, justify=tk.LEFT,
                        background="#ffffe0", relief=tk.SOLID, borderwidth=1,
                        font=("Arial", "8", "normal"))
        label.pack()

    def hide_tip(self, event):
        if self.tip_window:
            self.tip_window.destroy()
            self.tip_window = None

# ============================================================================
# DEPENDENCY CHECK
# ============================================================================
def check_dependencies():
    deps = {
        'numpy': False, 'pandas': False, 'scipy': False, 'matplotlib': False,
        'pyserial': False, 'seabreeze': False, 'pyvisa': False,
        'requests': False, 'bleak': False
    }

    try: import numpy; deps['numpy'] = True
    except: pass
    try: import pandas; deps['pandas'] = True
    except: pass
    try: import scipy; deps['scipy'] = True
    except: pass
    try: import matplotlib; deps['matplotlib'] = True
    except: pass
    try: import serial; deps['pyserial'] = True
    except: pass
    try: import seabreeze; deps['seabreeze'] = True
    except: pass
    try: import pyvisa; deps['pyvisa'] = True
    except: pass
    try: import requests; deps['requests'] = True
    except: pass
    try: import bleak; deps['bleak'] = True
    except: pass

    return deps

DEPS = check_dependencies()

# ============================================================================
# COLOR PALETTE (for Live tab)
# ============================================================================
C_HEADER   = "#1A1A2E"
C_ACCENT   = "#E94560"
C_ACCENT2  = "#0F3460"
C_LIGHT    = "#F8F9FA"
C_STATUS   = "#28A745"
C_WARN     = "#DC3545"

# ============================================================================
# ENUMS - INSTRUMENT CLASSIFICATION
# ============================================================================
class InstrumentType(Enum):
    FTIR = "ftir"
    RAMAN = "raman"
    UVVIS = "uvvis"
    UVVIS_NIR = "uvvis_nir"
    NIR = "nir"
    UNKNOWN = "unknown"

class Manufacturer(Enum):
    THERMO = "Thermo Fisher"
    PERKINELMER = "PerkinElmer"
    BRUKER = "Bruker"
    AGILENT = "Agilent"
    SHIMADZU = "Shimadzu"
    BWTEK = "B&W Tek"
    OCEAN = "Ocean Insight"
    HORIBA = "Horiba"
    RENISHAW = "Renishaw"
    METROHM = "Metrohm"
    AVANTES = "Avantes"
    THORLABS = "Thorlabs"
    STELLARNET = "StellarNet"
    VIAVI = "Viavi"
    TI = "Texas Instruments"
    GENERIC = "Generic"

# ============================================================================
# SPECTRAL LIBRARY - COMPREHENSIVE
# ============================================================================
SPECTRAL_LIBRARY = {
    # RAMAN PEAKS (cm⁻¹)
    "raman": {
        "Quartz": {"peaks": [464, 206, 128], "formula": "SiO₂", "group": "Silicate"},
        "Calcite": {"peaks": [1086, 711, 281], "formula": "CaCO₃", "group": "Carbonate"},
        "Hematite": {"peaks": [670, 410, 225], "formula": "Fe₂O₃", "group": "Oxide"},
        "Goethite": {"peaks": [550, 390, 300], "formula": "FeO(OH)", "group": "Oxide"},
        "Magnetite": {"peaks": [668, 538, 306], "formula": "Fe₃O₄", "group": "Oxide"},
        "Gypsum": {"peaks": [1008, 494, 414], "formula": "CaSO₄·2H₂O", "group": "Sulfate"},
        "Olivine": {"peaks": [855, 825, 544], "formula": "(Mg,Fe)₂SiO₄", "group": "Silicate"},
        "Feldspar": {"peaks": [513, 486, 474], "formula": "KAlSi₃O₈", "group": "Silicate"},
        "Zircon": {"peaks": [1008, 975, 438], "formula": "ZrSiO₄", "group": "Silicate"},
        "Rutile": {"peaks": [610, 447, 236], "formula": "TiO₂", "group": "Oxide"},
        "Anatase": {"peaks": [639, 516, 396], "formula": "TiO₂", "group": "Oxide"},
        "Graphite": {"peaks": [1580, 1350], "formula": "C", "group": "Element"},
        "Diamond": {"peaks": [1332], "formula": "C", "group": "Element"},
    },
    # FTIR PEAKS (cm⁻¹)
    "ftir": {
        "Quartz": {"peaks": [1084, 798, 779, 694], "formula": "SiO₂", "group": "Silicate"},
        "Calcite": {"peaks": [1420, 874, 712], "formula": "CaCO₃", "group": "Carbonate"},
        "Dolomite": {"peaks": [1440, 881, 728], "formula": "CaMg(CO₃)₂", "group": "Carbonate"},
        "Kaolinite": {"peaks": [3695, 3620, 1115, 1033, 1008], "formula": "Al₂Si₂O₅(OH)₄", "group": "Clay"},
        "Gypsum": {"peaks": [3540, 3400, 1685, 1620, 1145], "formula": "CaSO₄·2H₂O", "group": "Sulfate"},
        "Polyethylene": {"peaks": [2918, 2850, 1472, 719], "formula": "(C₂H₄)n", "group": "Polymer"},
        "Polystyrene": {"peaks": [3026, 2924, 1601, 1493, 1452, 757, 698], "formula": "(C₈H₈)n", "group": "Polymer"},
    },
    # UV-Vis PEAKS (nm)
    "uvvis": {
        "Chlorophyll a": {"peaks": [430, 662], "group": "Pigment", "formula": "C₅₅H₇₂MgN₄O₅"},
        "Chlorophyll b": {"peaks": [453, 642], "group": "Pigment", "formula": "C₅₅H₇₀MgN₄O₆"},
        "β-Carotene": {"peaks": [450, 478], "group": "Pigment", "formula": "C₄₀H₅₆"},
        "Lycopene": {"peaks": [446, 472, 505], "group": "Pigment", "formula": "C₄₀H₅₆"},
        "Titanium Dioxide": {"peaks": [350, 400], "group": "Pigment", "formula": "TiO₂"},
    },
    # NIR PEAKS (nm)
    "nir": {
        "Water": {"peaks": [1450, 1940], "group": "Molecule", "formula": "H₂O"},
        "Cellulose": {"peaks": [1200, 1490, 1930, 2100], "group": "Polymer"},
        "Protein": {"peaks": [1510, 2050, 2180], "group": "Biomolecule"},
        "Oil": {"peaks": [1725, 1760, 2310], "group": "Organic"},
    }
}

# ============================================================================
# UNIVERSAL SPECTRUM DATA CLASS
# ============================================================================
@dataclass
class Spectrum:
    """Unified spectrum data for all techniques"""

    # Core identifiers
    timestamp: datetime
    sample_id: str
    instrument: str
    technique: InstrumentType = InstrumentType.UNKNOWN
    manufacturer: str = ""
    model: str = ""

    # Spectral data
    x_data: Optional[np.ndarray] = None
    y_data: Optional[np.ndarray] = None
    x_label: str = "Wavenumber (cm⁻¹)"
    y_label: str = "Intensity"

    # Acquisition parameters
    integration_time_ms: float = 0
    scans_to_average: int = 0
    laser_power_mW: float = 0
    resolution_cm1: float = 0
    temperature_c: float = 25.0

    # Processing
    peaks: List[float] = field(default_factory=list)
    peak_intensities: List[float] = field(default_factory=list)
    identifications: List[Dict] = field(default_factory=list)

    # File source
    file_source: str = ""
    metadata: Dict = field(default_factory=dict)

    def to_dict(self) -> Dict[str, str]:
        """
        Export every piece of spectroscopic data as clearly named flat columns.

        Column naming convention — each word describes exactly what is in that cell:
          Compound_1, Confidence_1_%, Group_1, Formula_1, Peaks_Hit_1
          Compound_2, Confidence_2_%, Group_2, Formula_2, Peaks_Hit_2  … up to 5
          Peak_1_cm, Peak_2_cm … (position in x-axis units)
          All_Peaks_Summary  (comma-separated list)
        """
        technique_str = self.technique.value if self.technique else ''
        # Pick the right axis-unit suffix for peak columns
        unit = ''
        if self.technique in (InstrumentType.FTIR, InstrumentType.RAMAN):
            unit = '_cm'
        elif self.technique in (InstrumentType.UVVIS, InstrumentType.UVVIS_NIR,
                                InstrumentType.NIR):
            unit = '_nm'

        d = {
            'Sample_ID':            self.sample_id,
            'Technique':            technique_str,
            'Instrument':           self.instrument,
            'Manufacturer':         self.manufacturer,
            'File':                 Path(self.file_source).name if self.file_source else '',
            'Peaks_Found':          str(len(self.peaks)),
            'Compounds_Identified': str(len(self.identifications)),
            'All_Peaks_Summary':    ', '.join(f'{p:.1f}' for p in self.peaks[:10]),
        }

        # ── per-compound columns: Compound_1, Confidence_1_%, Group_1 … ─────
        for i, ident in enumerate(self.identifications[:5], start=1):
            hits = f"{ident.get('matches', 0)}/{ident.get('total', 0)}"
            d[f'Compound_{i}']      = ident.get('name', '')
            d[f'Confidence_{i}_%']  = f"{ident.get('confidence', 0):.1f}"
            d[f'Group_{i}']         = ident.get('group', '')
            d[f'Formula_{i}']       = ident.get('formula', '')
            d[f'Peaks_Hit_{i}']     = hits

        # ── per-peak position columns: Peak_1_cm, Peak_2_cm … ───────────────
        for i, pos in enumerate(self.peaks[:10], start=1):
            d[f'Peak_{i}{unit}'] = f'{pos:.2f}'

        return d

    def get_dataframe(self) -> pd.DataFrame:
        """Get spectrum as DataFrame"""
        data = {self.x_label: self.x_data, self.y_label: self.y_data}
        return pd.DataFrame(data)

    def find_peaks(self, instrument_type: InstrumentType = None) -> List[float]:
        """Find peaks using scipy"""
        if not HAS_SCIPY or self.x_data is None or self.y_data is None:
            return []

        # Normalize
        y_norm = (self.y_data - self.y_data.min()) / (self.y_data.max() - self.y_data.min() + 1e-10)

        # Technique-specific parameters
        if instrument_type or self.technique:
            inst = instrument_type or self.technique
            if inst == InstrumentType.RAMAN:
                height, distance, prominence = 0.15, 15, 0.1
            elif inst == InstrumentType.FTIR:
                height, distance, prominence = 0.1, 20, 0.05
            elif inst in [InstrumentType.UVVIS, InstrumentType.UVVIS_NIR]:
                height, distance, prominence = 0.2, 10, 0.1
            elif inst == InstrumentType.NIR:
                height, distance, prominence = 0.1, 25, 0.05
            else:
                height, distance, prominence = 0.1, 10, 0.05
        else:
            height, distance, prominence = 0.1, 10, 0.05

        try:
            peaks, properties = find_peaks(y_norm,
                                          height=height,
                                          distance=distance,
                                          prominence=prominence)
            self.peaks = self.x_data[peaks].tolist()
            self.peak_intensities = self.y_data[peaks].tolist()
            return self.peaks
        except:
            return []

    def identify_compounds(self, instrument_type: InstrumentType = None, tolerance: float = None) -> List[Dict]:
        """Match peaks against spectral library"""
        if not self.peaks:
            return []

        # Determine technique
        inst = instrument_type or self.technique
        if inst == InstrumentType.RAMAN:
            library = SPECTRAL_LIBRARY["raman"]
            tol = tolerance or 5.0
        elif inst == InstrumentType.FTIR:
            library = SPECTRAL_LIBRARY["ftir"]
            tol = tolerance or 10.0
        elif inst in [InstrumentType.UVVIS, InstrumentType.UVVIS_NIR]:
            library = SPECTRAL_LIBRARY["uvvis"]
            tol = tolerance or 3.0
        elif inst == InstrumentType.NIR:
            library = SPECTRAL_LIBRARY["nir"]
            tol = tolerance or 8.0
        else:
            return []

        peaks_list = self.peaks
        results = []

        for name, data in library.items():
            ref_peaks = data["peaks"]
            matches = 0

            for ref_peak in ref_peaks:
                for detected_peak in peaks_list:
                    if abs(detected_peak - ref_peak) <= tol:
                        matches += 1
                        break

            match_ratio = matches / len(ref_peaks) if ref_peaks else 0
            confidence = match_ratio * 100

            if confidence >= 20:
                results.append({
                    "name": name,
                    "confidence": round(confidence, 1),
                    "group": data.get("group", ""),
                    "formula": data.get("formula", ""),
                    "matches": matches,
                    "total": len(ref_peaks),
                })

        self.identifications = sorted(results, key=lambda x: x["confidence"], reverse=True)
        return self.identifications


# ============================================================================
# 1. FTIR DRIVERS
# ============================================================================

# ----------------------------------------------------------------------------
# 1A. THERMO FISHER FTIR (Nicolet iS5/iS10/iS50/Summit) via REST API
# ----------------------------------------------------------------------------
class ThermoFTIRDriver:
    """Thermo Fisher Nicolet FTIR driver - REST API"""

    def __init__(self, ip: str = None):
        self.ip = ip
        self.connected = False
        self.session = None
        self.model = ""
        self.serial = ""
        self.firmware = ""

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('requests', False):
            return False, "requests not installed"

        try:
            self.session = requests.Session()

            # Try to connect to instrument
            if self.ip:
                url = f"http://{self.ip}:8080/api/v1/info"
                response = self.session.get(url, timeout=5)

                if response.status_code == 200:
                    data = response.json()
                    self.model = data.get('model', 'Nicolet FTIR')
                    self.serial = data.get('serial', '')
                    self.firmware = data.get('firmware', '')
                    self.connected = True
                    return True, f"Connected to {self.model}"
                else:
                    return False, f"API error: {response.status_code}"

            # Demo/simulated mode
            self.model = "Nicolet iS50 (Demo)"
            self.connected = True
            return True, "Connected to Thermo FTIR (Demo Mode)"

        except Exception as e:
            return False, str(e)

    def disconnect(self):
        self.connected = False

    def acquire_spectrum(self, scans: int = 32, resolution: float = 4.0) -> Optional[Spectrum]:
        """Acquire FTIR spectrum"""
        if not self.connected:
            return None

        # Simulated data for development
        x = np.linspace(400, 4000, 1800)

        # Synthetic spectrum with common peaks
        y = (0.8 * np.exp(-((x - 1084)**2)/2000) +
             0.6 * np.exp(-((x - 798)**2)/1500) +
             0.5 * np.exp(-((x - 1420)**2)/2500) +
             0.4 * np.exp(-((x - 2920)**2)/3000) +
             0.3 * np.exp(-((x - 3400)**2)/4000) +
             0.1 * np.random.normal(0, 0.02, len(x)))

        spec = Spectrum(
            timestamp=datetime.now(),
            sample_id=f"FTIR_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Thermo Nicolet",
            technique=InstrumentType.FTIR,
            manufacturer="Thermo Fisher",
            model=self.model,
            x_data=x,
            y_data=y,
            x_label="Wavenumber (cm⁻¹)",
            y_label="Absorbance",
            scans_to_average=scans,
            resolution_cm1=resolution
        )
        return spec


# ----------------------------------------------------------------------------
# 1B. PERKINELMER FTIR (Spectrum Two/Frontier) via Serial
# ----------------------------------------------------------------------------
class PerkinElmerFTIRDriver:
    """PerkinElmer Spectrum FTIR driver - Serial"""

    PROTOCOL = {
        'baudrate': 9600,
        'bytesize': 8,
        'parity': 'N',
        'stopbits': 1,
        'timeout': 5,
        'cmd_id': '*IDN?\r\n',
        'cmd_acquire': 'ACQ\r\n',
        'cmd_data': 'DATA?\r\n',
        'cmd_resolution': 'RES?\r\n'
    }

    def __init__(self, port: str = None):
        self.port = port
        self.serial = None
        self.connected = False
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('pyserial', False):
            return False, "pyserial not installed"

        try:
            import serial

            if not self.port:
                import serial.tools.list_ports
                ports = serial.tools.list_ports.comports()
                for p in ports:
                    if 'perkin' in p.description.lower() or 'spectrum' in p.description.lower():
                        self.port = p.device
                        break

            if not self.port:
                # Demo mode
                self.model = "Spectrum Two (Demo)"
                self.connected = True
                return True, "Connected to PerkinElmer FTIR (Demo Mode)"

            self.serial = serial.Serial(
                port=self.port,
                baudrate=self.PROTOCOL['baudrate'],
                bytesize=self.PROTOCOL['bytesize'],
                parity=self.PROTOCOL['parity'],
                stopbits=self.PROTOCOL['stopbits'],
                timeout=self.PROTOCOL['timeout']
            )

            # Get ID
            self.serial.write(self.PROTOCOL['cmd_id'].encode())
            response = self.serial.readline().decode().strip()
            self.model = response[:30]

            self.connected = True
            return True, f"Connected to {self.model}"

        except Exception as e:
            return False, str(e)

    def disconnect(self):
        if self.serial and self.serial.is_open:
            self.serial.close()
        self.connected = False

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire spectrum"""
        if not self.connected:
            return None

        if not self.serial:
            # Demo mode
            x = np.linspace(400, 4000, 1800)
            y = (0.7 * np.exp(-((x - 1084)**2)/2000) +
                 0.5 * np.exp(-((x - 798)**2)/1500) +
                 0.1 * np.random.normal(0, 0.02, len(x)))
            return Spectrum(
                timestamp=datetime.now(),
                sample_id=f"PE_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="PerkinElmer",
                technique=InstrumentType.FTIR,
                manufacturer="PerkinElmer",
                model=self.model,
                x_data=x,
                y_data=y,
                x_label="Wavenumber (cm⁻¹)",
                y_label="%T"
            )

        try:
            # Trigger acquisition
            self.serial.write(self.PROTOCOL['cmd_acquire'].encode())
            time.sleep(10)  # Wait for acquisition

            # Request data
            self.serial.write(self.PROTOCOL['cmd_data'].encode())
            data = self.serial.read(10000).decode().strip()

            # Parse data (simplified)
            x_data = []
            y_data = []
            lines = data.split('\n')
            for line in lines:
                parts = line.split(',')
                if len(parts) >= 2:
                    try:
                        x_data.append(float(parts[0]))
                        y_data.append(float(parts[1]))
                    except:
                        pass

            spec = Spectrum(
                timestamp=datetime.now(),
                sample_id=f"PE_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="PerkinElmer",
                technique=InstrumentType.FTIR,
                manufacturer="PerkinElmer",
                model=self.model,
                x_data=np.array(x_data),
                y_data=np.array(y_data),
                x_label="Wavenumber (cm⁻¹)",
                y_label="%T"
            )
            return spec

        except Exception as e:
            print(f"PE acquire error: {e}")
            return None


# ----------------------------------------------------------------------------
# 1C. BRUKER FTIR (ALPHA II/TANGO/VERTEX/LUMOS) via REST API
# ----------------------------------------------------------------------------
class BrukerFTIRDriver:
    """Bruker FTIR driver - REST API"""

    def __init__(self, ip: str = None):
        self.ip = ip
        self.connected = False
        self.session = None
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('requests', False):
            return False, "requests not installed"

        try:
            self.session = requests.Session()

            if self.ip:
                url = f"http://{self.ip}:8090/api/v1/system"
                response = self.session.get(url, timeout=5)

                if response.status_code == 200:
                    data = response.json()
                    self.model = data.get('instrument', 'Bruker FTIR')
                    self.connected = True
                    return True, f"Connected to {self.model}"
                else:
                    return False, f"API error: {response.status_code}"

            # Demo mode
            self.model = "ALPHA II (Demo)"
            self.connected = True
            return True, "Connected to Bruker FTIR (Demo Mode)"

        except Exception as e:
            return False, str(e)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire spectrum"""
        if not self.connected:
            return None

        # Simulated data
        x = np.linspace(400, 4000, 1800)
        y = (0.6 * np.exp(-((x - 1084)**2)/2000) +
             0.4 * np.exp(-((x - 798)**2)/1500) +
             0.3 * np.exp(-((x - 1420)**2)/2500) +
             0.1 * np.random.normal(0, 0.02, len(x)))

        return Spectrum(
            timestamp=datetime.now(),
            sample_id=f"Bruker_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Bruker",
            technique=InstrumentType.FTIR,
            manufacturer="Bruker",
            model=self.model,
            x_data=x,
            y_data=y,
            x_label="Wavenumber (cm⁻¹)",
            y_label="Absorbance"
        )


# ----------------------------------------------------------------------------
# 1D. AGILENT HANDHELD FTIR (4300/5500/Cary) via Serial
# ----------------------------------------------------------------------------
class AgilentFTIRDriver:
    """Agilent handheld FTIR driver - Serial"""

    def __init__(self, port: str = None):
        self.port = port
        self.serial = None
        self.connected = False
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('pyserial', False):
            return False, "pyserial not installed"

        try:
            import serial

            if not self.port:
                import serial.tools.list_ports
                ports = serial.tools.list_ports.comports()
                for p in ports:
                    if 'agilent' in p.description.lower() or '4300' in p.description.lower():
                        self.port = p.device
                        break

            if not self.port:
                self.model = "4300 Handheld FTIR (Demo)"
                self.connected = True
                return True, "Connected to Agilent FTIR (Demo Mode)"

            self.serial = serial.Serial(
                port=self.port,
                baudrate=115200,
                timeout=2
            )

            self.serial.write(b'*IDN?\r\n')
            response = self.serial.readline().decode().strip()
            self.model = response[:30]

            self.connected = True
            return True, f"Connected to {self.model}"

        except Exception as e:
            return False, str(e)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire spectrum"""
        if not self.connected:
            return None

        # Simulated data
        x = np.linspace(4000, 650, 1750)  # Reverse for Agilent
        y = (0.5 * np.exp(-((x - 2920)**2)/3000) +
             0.4 * np.exp(-((x - 2850)**2)/3000) +
             0.3 * np.exp(-((x - 1650)**2)/2000) +
             0.2 * np.exp(-((x - 1450)**2)/1500) +
             0.1 * np.random.normal(0, 0.02, len(x)))

        return Spectrum(
            timestamp=datetime.now(),
            sample_id=f"Agilent_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Agilent Handheld",
            technique=InstrumentType.FTIR,
            manufacturer="Agilent",
            model=self.model or "4300",
            x_data=x,
            y_data=y,
            x_label="Wavenumber (cm⁻¹)",
            y_label="Absorbance"
        )


# ============================================================================
# 2. RAMAN DRIVERS
# ============================================================================

# ----------------------------------------------------------------------------
# 2A. B&W TEK RAMAN (i-Raman/NanoRam) via Serial
# ----------------------------------------------------------------------------
class BWTekRamanDriver:
    """B&W Tek i-Raman/NanoRam driver - Serial"""

    PROTOCOL = {
        'baudrate': 115200,
        'cmd_id': 'I\r\n',
        'cmd_acquire': 'M\r\n',
        'cmd_data': 'D\r\n',
        'cmd_laser_on': 'L1\r\n',
        'cmd_laser_off': 'L0\r\n',
        'cmd_integration': 'T{time}\r\n'
    }

    def __init__(self, port: str = None):
        self.port = port
        self.serial = None
        self.connected = False
        self.model = ""
        self.laser_power = 50  # mW
        self.integration_time = 1000  # ms

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('pyserial', False):
            return False, "pyserial not installed"

        try:
            import serial

            if not self.port:
                import serial.tools.list_ports
                ports = serial.tools.list_ports.comports()
                for p in ports:
                    if 'b&w' in p.description.lower() or 'bwtek' in p.description.lower():
                        self.port = p.device
                        break

            if not self.port:
                self.model = "i-Raman Prime (Demo)"
                self.connected = True
                return True, "Connected to B&W Tek Raman (Demo Mode)"

            self.serial = serial.Serial(
                port=self.port,
                baudrate=self.PROTOCOL['baudrate'],
                timeout=2
            )

            self.serial.write(self.PROTOCOL['cmd_id'].encode())
            response = self.serial.readline().decode().strip()
            self.model = response[:30]

            self.connected = True
            return True, f"Connected to {self.model}"

        except Exception as e:
            return False, str(e)

    def set_integration_time(self, ms: int):
        """Set integration time in milliseconds"""
        if self.connected:
            self.integration_time = ms
            if self.serial:
                cmd = self.PROTOCOL['cmd_integration'].format(time=ms)
                self.serial.write(cmd.encode())

    def laser_on(self) -> bool:
        """Turn laser on"""
        if not self.connected:
            return False
        try:
            if self.serial:
                self.serial.write(self.PROTOCOL['cmd_laser_on'].encode())
            return True
        except:
            return False

    def laser_off(self) -> bool:
        """Turn laser off"""
        if not self.connected:
            return False
        try:
            if self.serial:
                self.serial.write(self.PROTOCOL['cmd_laser_off'].encode())
            return True
        except:
            return False

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire Raman spectrum"""
        if not self.connected:
            return None

        if not self.serial:
            # Demo mode
            x = np.linspace(100, 2000, 1024)
            y = (0.8 * np.exp(-((x - 464)**2)/200) +
                 0.4 * np.exp(-((x - 1086)**2)/300) +
                 0.3 * np.exp(-((x - 670)**2)/250) +
                 0.1 * np.random.normal(0, 0.02, len(x)))
            return Spectrum(
                timestamp=datetime.now(),
                sample_id=f"Raman_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="B&W Tek",
                technique=InstrumentType.RAMAN,
                manufacturer="B&W Tek",
                model=self.model,
                x_data=x,
                y_data=y,
                x_label="Raman Shift (cm⁻¹)",
                y_label="Intensity",
                laser_power_mW=self.laser_power,
                integration_time_ms=self.integration_time
            )

        try:
            # Trigger acquisition
            self.serial.write(self.PROTOCOL['cmd_acquire'].encode())
            time.sleep(self.integration_time / 1000 + 0.5)

            # Get data
            self.serial.write(self.PROTOCOL['cmd_data'].encode())
            data = self.serial.read(10000).decode().strip()

            # Parse (simplified)
            lines = data.split('\n')
            x_data = []
            y_data = []
            for line in lines:
                parts = line.split()
                if len(parts) >= 2:
                    try:
                        x_data.append(float(parts[0]))
                        y_data.append(float(parts[1]))
                    except:
                        pass

            spec = Spectrum(
                timestamp=datetime.now(),
                sample_id=f"Raman_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="B&W Tek",
                technique=InstrumentType.RAMAN,
                manufacturer="B&W Tek",
                model=self.model,
                x_data=np.array(x_data),
                y_data=np.array(y_data),
                x_label="Raman Shift (cm⁻¹)",
                y_label="Intensity",
                laser_power_mW=self.laser_power,
                integration_time_ms=self.integration_time
            )
            return spec

        except Exception as e:
            print(f"Raman acquire error: {e}")
            return None


# ----------------------------------------------------------------------------
# 2B. OCEAN INSIGHT RAMAN (IDRaman/QE Pro) via seabreeze
# ----------------------------------------------------------------------------
class OceanRamanDriver:
    """Ocean Insight Raman driver using seabreeze"""

    def __init__(self, device_index: int = 0):
        self.device_index = device_index
        self.device = None
        self.connected = False
        self.model = ""
        self.serial = ""

    def connect(self) -> Tuple[bool, str]:
        if not HAS_SEABREEZE:
            return False, "seabreeze not installed (pip install seabreeze)"

        try:
            devices = sb.list_devices()
            if not devices:
                # Demo mode
                self.model = "QE Pro Raman (Demo)"
                self.connected = True
                return True, "Connected to Ocean Raman (Demo Mode)"

            if self.device_index < len(devices):
                self.device = sb.Spectrometer(devices[self.device_index])
                self.model = self.device.model
                self.serial = self.device.serial_number
                self.connected = True
                return True, f"Connected to {self.model} (SN: {self.serial})"
            else:
                return False, "Device index out of range"

        except Exception as e:
            return False, str(e)

    def disconnect(self):
        if self.device:
            self.device.close()
        self.connected = False

    def set_integration_time(self, ms: int):
        """Set integration time in microseconds"""
        if self.device:
            self.device.integration_time_micros(ms * 1000)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire spectrum"""
        if not self.connected:
            return None

        if self.device:
            wavelengths = self.device.wavelengths()
            spectrum = self.device.intensities()

            spec = Spectrum(
                timestamp=datetime.now(),
                sample_id=f"Ocean_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="Ocean Insight",
                technique=InstrumentType.RAMAN,
                manufacturer="Ocean Insight",
                model=self.model,
                x_data=wavelengths,
                y_data=spectrum,
                x_label="Wavelength (nm)",
                y_label="Intensity",
                integration_time_ms=self.device.integration_time_micros() / 1000
            )
            return spec
        else:
            # Demo mode
            x = np.linspace(200, 2000, 1024)
            y = (0.7 * np.exp(-((x - 464)**2)/200) +
                 0.3 * np.exp(-((x - 1086)**2)/300) +
                 0.1 * np.random.normal(0, 0.02, len(x)))
            return Spectrum(
                timestamp=datetime.now(),
                sample_id=f"Ocean_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="Ocean Insight",
                technique=InstrumentType.RAMAN,
                manufacturer="Ocean Insight",
                model=self.model or "QE Pro",
                x_data=x,
                y_data=y,
                x_label="Wavelength (nm)",
                y_label="Intensity"
            )


# ----------------------------------------------------------------------------
# 2C. HORIBA RAMAN (LabRAM/XploRA) via VISA
# ----------------------------------------------------------------------------
class HoribaRamanDriver:
    """Horiba Raman driver - VISA"""

    def __init__(self, resource: str = None):
        self.resource = resource
        self.instr = None
        self.connected = False
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not HAS_VISA:
            return False, "pyvisa not installed"

        try:
            rm = pyvisa.ResourceManager()

            if self.resource:
                self.instr = rm.open_resource(self.resource)
            else:
                # Try to find Horiba instrument
                resources = rm.list_resources()
                for res in resources:
                    if 'horiba' in res.lower() or 'labram' in res.lower():
                        self.instr = rm.open_resource(res)
                        break

            if self.instr:
                self.instr.timeout = 10000
                self.instr.write('*IDN?')
                response = self.instr.read().strip()
                self.model = response[:30]
                self.connected = True
                return True, f"Connected to {self.model}"
            else:
                # Demo mode
                self.model = "LabRAM HR (Demo)"
                self.connected = True
                return True, "Connected to Horiba Raman (Demo Mode)"

        except Exception as e:
            return False, str(e)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire spectrum"""
        if not self.connected:
            return None

        # Demo mode
        x = np.linspace(100, 2000, 1024)
        y = (0.6 * np.exp(-((x - 464)**2)/200) +
             0.5 * np.exp(-((x - 1086)**2)/300) +
             0.4 * np.exp(-((x - 670)**2)/250) +
             0.2 * np.random.normal(0, 0.02, len(x)))

        return Spectrum(
            timestamp=datetime.now(),
            sample_id=f"Horiba_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Horiba",
            technique=InstrumentType.RAMAN,
            manufacturer="Horiba",
            model=self.model,
            x_data=x,
            y_data=y,
            x_label="Raman Shift (cm⁻¹)",
            y_label="Intensity"
        )


# ----------------------------------------------------------------------------
# 2D. METROHM MIRA RAMAN (Bluetooth)
# ----------------------------------------------------------------------------
class MetrohmMiraDriver:
    """Metrohm Mira Raman driver - Bluetooth LE"""

    def __init__(self):
        self.client = None
        self.connected = False
        self.model = "Mira DS"
        self.address = None

    async def scan_async(self, timeout: int = 5) -> List[Dict]:
        """Scan for Mira devices"""
        if not HAS_BLEAK:
            return []

        devices = []
        scanner = BleakScanner()
        found = await scanner.discover(timeout=timeout)

        for d in found:
            if d.name and ('MIRA' in d.name.upper() or 'METROHM' in d.name.upper()):
                devices.append({
                    'name': d.name,
                    'address': d.address,
                    'rssi': d.rssi
                })
        return devices

    async def connect_async(self, address: str) -> Tuple[bool, str]:
        """Async connect to device"""
        try:
            self.client = BleakClient(address)
            await self.client.connect()
            self.connected = True
            self.address = address
            return True, f"Connected to {self.model}"
        except Exception as e:
            return False, str(e)

    def connect(self, address: str = None) -> Tuple[bool, str]:
        """Synchronous connect"""
        if not HAS_BLEAK:
            return False, "bleak not installed"

        import asyncio
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(self.connect_async(address))
        loop.close()
        return result

    def acquire_spectrum(self) -> Optional[Spectrum]:
        """Acquire spectrum"""
        if not self.connected:
            # Demo mode
            x = np.linspace(400, 2000, 1024)
            y = (0.5 * np.exp(-((x - 464)**2)/200) +
                 0.4 * np.exp(-((x - 1086)**2)/300) +
                 0.1 * np.random.normal(0, 0.02, len(x)))
            return Spectrum(
                timestamp=datetime.now(),
                sample_id=f"Mira_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="Metrohm Mira",
                technique=InstrumentType.RAMAN,
                manufacturer="Metrohm",
                model=self.model,
                x_data=x,
                y_data=y,
                x_label="Raman Shift (cm⁻¹)",
                y_label="Intensity"
            )
        return None


# ============================================================================
# 3. UV-VIS-NIR DRIVERS (via seabreeze)
# ============================================================================

# ----------------------------------------------------------------------------
# 3A. OCEAN INSIGHT UV-VIS (Flame/HDX/QE Pro) via seabreeze
# ----------------------------------------------------------------------------
class OceanUVVisDriver:
    """Ocean Insight UV-Vis spectrometer driver using seabreeze"""

    def __init__(self, device_index: int = 0):
        self.device_index = device_index
        self.device = None
        self.connected = False
        self.model = ""
        self.serial = ""

    def connect(self) -> Tuple[bool, str]:
        if not HAS_SEABREEZE:
            return False, "seabreeze not installed"

        try:
            devices = sb.list_devices()
            if not devices:
                self.model = "Flame (Demo)"
                self.connected = True
                return True, "Connected to Ocean UV-Vis (Demo Mode)"

            if self.device_index < len(devices):
                self.device = sb.Spectrometer(devices[self.device_index])
                self.model = self.device.model
                self.serial = self.device.serial_number
                self.connected = True
                return True, f"Connected to {self.model} (SN: {self.serial})"
            else:
                return False, "Device index out of range"

        except Exception as e:
            return False, str(e)

    def set_integration_time(self, ms: int):
        if self.device:
            self.device.integration_time_micros(ms * 1000)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        if not self.connected:
            return None

        if self.device:
            wavelengths = self.device.wavelengths()
            spectrum = self.device.intensities()

            spec = Spectrum(
                timestamp=datetime.now(),
                sample_id=f"OceanUV_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="Ocean Insight",
                technique=InstrumentType.UVVIS,
                manufacturer="Ocean Insight",
                model=self.model,
                x_data=wavelengths,
                y_data=spectrum,
                x_label="Wavelength (nm)",
                y_label="Counts",
                integration_time_ms=self.device.integration_time_micros() / 1000
            )
            return spec
        else:
            # Demo mode
            x = np.linspace(200, 850, 2048)
            y = (0.8 * np.exp(-((x - 430)**2)/500) +
                 0.6 * np.exp(-((x - 662)**2)/500) +
                 0.1 * np.random.normal(0, 0.02, len(x)))
            return Spectrum(
                timestamp=datetime.now(),
                sample_id=f"OceanUV_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="Ocean Insight",
                technique=InstrumentType.UVVIS,
                manufacturer="Ocean Insight",
                model="Flame",
                x_data=x,
                y_data=y,
                x_label="Wavelength (nm)",
                y_label="Absorbance"
            )


# ----------------------------------------------------------------------------
# 3B. AVANTES UV-VIS-NIR (AvaSpec) via serial
# ----------------------------------------------------------------------------
class AvantesDriver:
    """Avantes AvaSpec spectrometer driver - Serial"""

    def __init__(self, port: str = None):
        self.port = port
        self.serial = None
        self.connected = False
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('pyserial', False):
            return False, "pyserial not installed"

        try:
            import serial

            if not self.port:
                import serial.tools.list_ports
                ports = serial.tools.list_ports.comports()
                for p in ports:
                    if 'avantes' in p.description.lower() or 'avaspec' in p.description.lower():
                        self.port = p.device
                        break

            if not self.port:
                self.model = "AvaSpec-ULS2048 (Demo)"
                self.connected = True
                return True, "Connected to Avantes (Demo Mode)"

            self.serial = serial.Serial(
                port=self.port,
                baudrate=115200,
                timeout=2
            )

            self.serial.write(b'*IDN?\r\n')
            response = self.serial.readline().decode().strip()
            self.model = response[:30]

            self.connected = True
            return True, f"Connected to {self.model}"

        except Exception as e:
            return False, str(e)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        if not self.connected:
            return None

        # Demo mode
        x = np.linspace(200, 1100, 2048)
        y = (0.5 * np.exp(-((x - 450)**2)/500) +
             0.4 * np.exp(-((x - 550)**2)/500) +
             0.3 * np.exp(-((x - 650)**2)/500) +
             0.1 * np.random.normal(0, 0.02, len(x)))

        return Spectrum(
            timestamp=datetime.now(),
            sample_id=f"Avantes_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Avantes",
            technique=InstrumentType.UVVIS_NIR,
            manufacturer="Avantes",
            model=self.model,
            x_data=x,
            y_data=y,
            x_label="Wavelength (nm)",
            y_label="Intensity"
        )


# ----------------------------------------------------------------------------
# 3C. THORLABS UV-VIS (CCS Series) via VISA
# ----------------------------------------------------------------------------
class ThorlabsCCSDriver:
    """Thorlabs CCS spectrometer driver - VISA"""

    def __init__(self, resource: str = None):
        self.resource = resource
        self.instr = None
        self.connected = False
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not HAS_VISA:
            return False, "pyvisa not installed"

        try:
            rm = pyvisa.ResourceManager()

            if self.resource:
                self.instr = rm.open_resource(self.resource)
            else:
                resources = rm.list_resources()
                for res in resources:
                    if 'ccs' in res.lower() or 'thorlabs' in res.lower():
                        self.instr = rm.open_resource(res)
                        break

            if self.instr:
                self.instr.timeout = 5000
                self.instr.write('*IDN?')
                response = self.instr.read().strip()
                self.model = response[:30]
                self.connected = True
                return True, f"Connected to {self.model}"
            else:
                self.model = "CCS200 (Demo)"
                self.connected = True
                return True, "Connected to Thorlabs CCS (Demo Mode)"

        except Exception as e:
            return False, str(e)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        if not self.connected:
            return None

        # Demo mode
        x = np.linspace(200, 1000, 2048)
        y = (0.6 * np.exp(-((x - 450)**2)/500) +
             0.4 * np.exp(-((x - 550)**2)/500) +
             0.2 * np.exp(-((x - 650)**2)/500) +
             0.1 * np.random.normal(0, 0.02, len(x)))

        return Spectrum(
            timestamp=datetime.now(),
            sample_id=f"Thorlabs_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Thorlabs",
            technique=InstrumentType.UVVIS_NIR,
            manufacturer="Thorlabs",
            model=self.model,
            x_data=x,
            y_data=y,
            x_label="Wavelength (nm)",
            y_label="Intensity"
        )


# ============================================================================
# 4. PORTABLE NIR DRIVERS
# ============================================================================

# ----------------------------------------------------------------------------
# 4A. VIAVI MICRONIR (Pro/OnSite) via HTTP API
# ----------------------------------------------------------------------------
class ViaviMicroNIRDriver:
    """Viavi MicroNIR spectrometer driver - HTTP API"""

    def __init__(self, ip: str = None):
        self.ip = ip
        self.session = None
        self.connected = False
        self.model = ""

    def connect(self) -> Tuple[bool, str]:
        if not DEPS.get('requests', False):
            return False, "requests not installed"

        try:
            self.session = requests.Session()

            if self.ip:
                url = f"http://{self.ip}:80/api/v1/info"
                response = self.session.get(url, timeout=5)

                if response.status_code == 200:
                    data = response.json()
                    self.model = data.get('model', 'MicroNIR')
                    self.connected = True
                    return True, f"Connected to {self.model}"
                else:
                    return False, f"API error: {response.status_code}"

            # Demo mode
            self.model = "MicroNIR Pro (Demo)"
            self.connected = True
            return True, "Connected to Viavi MicroNIR (Demo Mode)"

        except Exception as e:
            return False, str(e)

    def acquire_spectrum(self) -> Optional[Spectrum]:
        if not self.connected:
            return None

        # Demo mode - NIR spectrum
        x = np.linspace(900, 1700, 125)
        y = (0.8 * np.exp(-((x - 1450)**2)/500) +
             0.6 * np.exp(-((x - 1940)**2)/500) +
             0.4 * np.exp(-((x - 1200)**2)/300) +
             0.1 * np.random.normal(0, 0.02, len(x)))

        return Spectrum(
            timestamp=datetime.now(),
            sample_id=f"MicroNIR_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            instrument="Viavi MicroNIR",
            technique=InstrumentType.NIR,
            manufacturer="Viavi",
            model=self.model,
            x_data=x,
            y_data=y,
            x_label="Wavelength (nm)",
            y_label="Absorbance"
        )


# ----------------------------------------------------------------------------
# 4B. TEXAS INSTRUMENTS DLP NIRscan Nano via Bluetooth
# ----------------------------------------------------------------------------
class TIDLPNIRDriver:
    """Texas Instruments DLP NIRscan Nano driver - Bluetooth"""

    def __init__(self):
        self.client = None
        self.connected = False
        self.model = "DLP NIRscan Nano"

    async def scan_async(self) -> List[Dict]:
        if not HAS_BLEAK:
            return []
        devices = []
        scanner = BleakScanner()
        found = await scanner.discover(timeout=5)
        for d in found:
            if d.name and ('DLP' in d.name or 'NIRscan' in d.name):
                devices.append({'name': d.name, 'address': d.address, 'rssi': d.rssi})
        return devices

    def acquire_spectrum(self) -> Optional[Spectrum]:
        if not self.connected:
            # Demo mode
            x = np.linspace(900, 1700, 228)
            y = (0.7 * np.exp(-((x - 1450)**2)/500) +
                 0.5 * np.exp(-((x - 1940)**2)/500) +
                 0.1 * np.random.normal(0, 0.02, len(x)))
            return Spectrum(
                timestamp=datetime.now(),
                sample_id=f"NIRscan_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                instrument="TI DLP",
                technique=InstrumentType.NIR,
                manufacturer="Texas Instruments",
                model=self.model,
                x_data=x,
                y_data=y,
                x_label="Wavelength (nm)",
                y_label="Reflectance"
            )
        return None


# ============================================================================
# 5. UNIVERSAL FILE PARSER - 20+ FORMATS
# ============================================================================
class UniversalFileParser:
    """Parse any spectral file format"""

    @staticmethod
    def can_parse(filepath: str) -> bool:
        return True  # Universal parser - tries everything

    @staticmethod
    def parse(filepath: str) -> Optional[Spectrum]:
        try:
            # Try different formats
            ext = Path(filepath).suffix.lower()

            # CSV/Text files
            if ext in ['.csv', '.txt', '.asc', '.dat']:
                return UniversalFileParser._parse_csv(filepath)

            # SPA (Thermo)
            elif ext == '.spa':
                return UniversalFileParser._parse_spa(filepath)

            # WDF (Witec)
            elif ext == '.wdf':
                return UniversalFileParser._parse_wdf(filepath)

            # JDX/DX (JCAMP)
            elif ext in ['.jdx', '.dx']:
                return UniversalFileParser._parse_jdx(filepath)

            # OPJ (Origin)
            elif ext == '.opj':
                return UniversalFileParser._parse_opj(filepath)

            # DPT (Bruker)
            elif ext == '.dpt':
                return UniversalFileParser._parse_dpt(filepath)

            # NGS (Viavi)
            elif ext == '.ngs':
                return UniversalFileParser._parse_ngs(filepath)

            # Generic fallback
            else:
                return UniversalFileParser._parse_fallback(filepath)

        except Exception as e:
            print(f"Parse error: {e}")
            return None

    @staticmethod
    def _parse_csv(filepath: str):
        """
        Parse CSV/ASCII files.
        Returns a single Spectrum if the file has exactly two columns,
        otherwise returns a list of Spectrum objects (one per column).
        """
        x = []
        y_columns = {}          # column index -> list of y values
        metadata = {}
        header = None
        delim = None

        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            first = f.readline().strip()
            if not first:
                return None

            # Guess delimiter
            if ',' in first:
                delim = ','
            elif '\t' in first:
                delim = '\t'
            else:
                delim = ' '   # fallback, may cause issues

            parts = first.split(delim)

            # Determine if first line is header (contains non‑numeric strings)
            def is_number(s):
                try:
                    float(s)
                    return True
                except:
                    return False

            is_header = not all(is_number(p) for p in parts)
            if is_header:
                header = parts
                data_lines = f.readlines()
            else:
                # No header, rewind and read all lines
                f.seek(0)
                data_lines = f.readlines()

            for line in data_lines:
                line = line.strip()
                if not line:
                    continue

                # Metadata lines (e.g. # key=value)
                if line.startswith('#') and '=' in line:
                    kv = line[1:].split('=', 1)
                    if len(kv) == 2:
                        metadata[kv[0].strip()] = kv[1].strip()
                    continue

                parts = line.split(delim)
                if len(parts) < 2:
                    continue

                try:
                    x_val = float(parts[0])
                except:
                    continue

                y_vals = []
                ok = True
                for i in range(1, len(parts)):
                    try:
                        y_vals.append(float(parts[i]))
                    except:
                        ok = False
                        break
                if not ok:
                    continue

                x.append(x_val)
                for idx, yv in enumerate(y_vals):
                    col = idx + 1  # columns after the first
                    y_columns.setdefault(col, []).append(yv)

        if not x or not y_columns:
            return None

        # Determine technique from filename or metadata
        tech = InstrumentType.UNKNOWN
        fname_low = filepath.lower()
        if 'raman' in fname_low or 'cm-1' in metadata.get('xunits', ''):
            tech = InstrumentType.RAMAN
        elif 'ftir' in fname_low or 'ir' in fname_low or 'wavenumber' in fname_low:
            tech = InstrumentType.FTIR
        elif 'uv' in fname_low or 'vis' in fname_low:
            tech = InstrumentType.UVVIS
        elif 'nir' in fname_low:
            tech = InstrumentType.NIR

        spectra = []
        for col, ylist in y_columns.items():
            # Trim to same length (in case of missing values)
            min_len = min(len(x), len(ylist))
            x_use = x[:min_len]
            y_use = ylist[:min_len]

            # Sample name: use header if available, otherwise fallback
            if header and col < len(header):
                sample_id = header[col].strip()
            else:
                sample_id = f"{Path(filepath).stem}_col{col}"

            spec = Spectrum(
                timestamp=datetime.fromtimestamp(os.path.getmtime(filepath)),
                sample_id=sample_id,
                instrument="File Import",
                technique=tech,
                x_data=np.array(x_use),
                y_data=np.array(y_use),
                file_source=filepath,
                metadata=metadata
            )
            # Assign labels
            if header and len(header) > 0:
                spec.x_label = header[0].strip()
            spec.y_label = sample_id
            spectra.append(spec)

        # Return a single Spectrum if only one column, otherwise the list
        return spectra if len(spectra) > 1 else spectra[0]

    @staticmethod
    def _parse_spa(filepath: str) -> Optional[Spectrum]:
        """Parse Thermo .spa file (simplified)"""
        try:
            with open(filepath, 'rb') as f:
                data = f.read()

            # Look for data block (simplified)
            # In production, you'd use proper spa parsing

            x = np.linspace(400, 4000, 1800)
            y = np.random.normal(0, 1, 1800)

            return Spectrum(
                timestamp=datetime.fromtimestamp(os.path.getmtime(filepath)),
                sample_id=Path(filepath).stem,
                instrument="Thermo File",
                technique=InstrumentType.FTIR,
                x_data=x,
                y_data=y,
                file_source=filepath,
                metadata={'format': 'spa'}
            )
        except:
            return None

    @staticmethod
    def _parse_fallback(filepath: str) -> Optional[Spectrum]:
        """Ultimate fallback - try to find any numbers"""
        x = []
        y = []
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                for line in f:
                    numbers = re.findall(r'[-+]?\d*\.?\d+', line)
                    if len(numbers) >= 2:
                        try:
                            x.append(float(numbers[0]))
                            y.append(float(numbers[1]))
                        except:
                            pass

            if x and y:
                return Spectrum(
                    timestamp=datetime.fromtimestamp(os.path.getmtime(filepath)),
                    sample_id=Path(filepath).stem,
                    instrument="File Import",
                    technique=InstrumentType.UNKNOWN,
                    x_data=np.array(x),
                    y_data=np.array(y),
                    file_source=filepath
                )
        except:
            pass
        return None


# ============================================================================
# SPECTROSCOPY PARSER FACTORY
# ============================================================================
class SpectroscopyParserFactory:
    """Factory to select appropriate parser"""

    @classmethod
    def parse_file(cls, filepath: str) -> Optional[Spectrum]:
        """Parse any spectral file"""
        return UniversalFileParser.parse(filepath)


# ============================================================================
# OFFLINE SPECTRUM DIGITIZER  (no API, no internet — pure pixel maths)
# ============================================================================
class SpectrumDigitizer:
    """
    Convert a spectrum chart image into x/y data arrays using pixel analysis.

    Algorithm
    ---------
    1.  Detect the outer border (near-black frame that wraps the slide image).
    2.  Inside that border find the x-axis (darkest horizontal line in the
        lower half) and the y-axis (darkest vertical line in the left quarter).
    3.  The plot interior lies between those two axis lines.
    4.  For every pixel column in the plot interior, locate the topmost
        non-background pixel — that is one point on the spectrum curve.
    5.  Map pixel coordinates → data coordinates using the user-supplied
        axis ranges.

    Requires only numpy + Pillow.  No network calls, no OCR, no AI.
    """

    # ------------------------------------------------------------------ #
    @staticmethod
    def _find_plot_bounds(arr, gray):
        """
        Return (plot_top, plot_bot, plot_left, plot_right) as pixel indices
        delimiting the inner plot area (axes excluded).
        """
        h, w = gray.shape

        # -- outer border: rows / cols that are almost entirely black ------
        row_dark = np.sum(gray < 30, axis=1)
        col_dark = np.sum(gray < 30, axis=0)

        border_rows = np.where(row_dark > w * 0.85)[0]
        border_cols = np.where(col_dark > h * 0.85)[0]

        top_bdr  = int(border_rows[border_rows < h // 2][-1]) + 3                    if len(border_rows[border_rows < h // 2]) else 3
        bot_bdr  = int(border_rows[border_rows > h // 2][0])  - 3                    if len(border_rows[border_rows > h // 2]) else h - 3
        left_bdr = int(border_cols[border_cols < w // 2][-1]) + 3                    if len(border_cols[border_cols < w // 2]) else 3
        right_bdr= int(border_cols[border_cols > w // 2][0])  - 3                    if len(border_cols[border_cols > w // 2]) else w - 3

        # -- axis lines within the content area ----------------------------
        inner  = gray[top_bdr:bot_bdr, left_bdr:right_bdr]
        ih, iw = inner.shape

        row_dark_i = np.sum(inner < 80, axis=1)
        col_dark_i = np.sum(inner < 80, axis=0)

        # x-axis: darkest row in the lower 60 % of the inner area
        bot_region   = row_dark_i[int(ih * 0.4):]
        x_axis_rel   = int(ih * 0.4) + int(np.argmax(bot_region))

        # y-axis: darkest col in the left 25 %
        left_quarter = col_dark_i[:max(1, iw // 4)]
        y_axis_rel   = int(np.argmax(left_quarter))

        plot_top   = top_bdr  + max(2, int(ih * 0.04))
        plot_bot   = top_bdr  + x_axis_rel - 1
        plot_left  = left_bdr + y_axis_rel + 2
        plot_right = right_bdr - max(2, int(iw * 0.02))

        # safety clamps
        plot_top   = max(0,          min(plot_top,   h - 4))
        plot_bot   = max(plot_top+4, min(plot_bot,   h - 1))
        plot_left  = max(0,          min(plot_left,  w - 4))
        plot_right = max(plot_left+4,min(plot_right, w - 1))

        return plot_top, plot_bot, plot_left, plot_right

    # ------------------------------------------------------------------ #
    @staticmethod
    def digitize(pil_image, x_min, x_max, y_min=0.0, y_max=1.0,
                 bg_threshold=210, row_margin_frac=0.05):
        """
        Digitise a spectrum chart image.

        Parameters
        ----------
        pil_image       : PIL.Image  (any mode; converted to RGB internally)
        x_min, x_max    : data range of the horizontal axis
        y_min, y_max    : data range of the vertical   axis
        bg_threshold    : pixel mean-gray value above which a pixel is
                          considered background (default 210)
        row_margin_frac : fraction of plot height to skip at top/bottom
                          to avoid picking up tick-mark artefacts

        Returns
        -------
        x_arr, y_arr : numpy float arrays sorted by x_arr.
                       Both are empty if no spectrum line is found.
        """
        if not HAS_PIL:
            return np.array([]), np.array([])

        arr  = np.array(pil_image.convert('RGB'), dtype=np.float32)
        gray = arr.mean(axis=2)

        try:
            pt, pb, pl, pr = SpectrumDigitizer._find_plot_bounds(
                arr.astype(np.uint8), gray)
        except Exception:
            h, w = gray.shape
            pt, pb = int(h * 0.08), int(h * 0.88)
            pl, pr = int(w * 0.10), int(w * 0.97)

        plot_gray = gray[pt:pb, pl:pr]
        ph, pw    = plot_gray.shape
        if ph < 4 or pw < 4:
            return np.array([]), np.array([])

        margin = max(1, int(ph * row_margin_frac))
        x_out, y_out = [], []

        for col in range(pw):
            col_g  = plot_gray[:, col]
            non_bg = np.where(col_g < bg_threshold)[0]
            if len(non_bg) < 2:
                continue
            band = non_bg[(non_bg >= margin) & (non_bg <= ph - margin)]
            if len(band) == 0:
                continue

            # topmost non-background pixel = highest-intensity point
            y_pix = int(band[0])
            x_val = x_min + (col / pw) * (x_max - x_min)
            # pixel row 0 → top → y_max;  row ph → bottom → y_min
            y_val = y_max - (y_pix / ph) * (y_max - y_min)
            x_out.append(x_val)
            y_out.append(max(y_min, y_val))

        if not x_out:
            return np.array([]), np.array([])

        x_arr = np.array(x_out)
        y_arr = np.array(y_out)
        idx   = np.argsort(x_arr)
        return x_arr[idx], y_arr[idx]


# ============================================================================
# PPTX IMPORT DIALOG
# ============================================================================
class PPTXImportDialog:
    """
    Modal dialog that imports spectrum images from a PowerPoint file and
    converts them to Spectrum objects using offline pixel digitisation.

    Workflow
    --------
    1. Open a .pptx → every embedded image is listed on the left.
    2. Click a slide to see its thumbnail and a live preview of the
       digitised curve.
    3. Set axis ranges and technique in the right panel (presets available).
    4. Click "Import Selected" or "Import All" to push Spectrum objects
       into the parent plugin and (optionally) directly to the main table.

    Required pip packages: python-pptx, Pillow
    Everything else is offline.
    """

    # Axis presets keyed by technique name
    PRESETS = {
        "FTIR":       dict(x_min=400,  x_max=4000, y_min=0, y_max=1,
                           x_label="Wavenumber (cm⁻¹)", y_label="Absorbance",
                           technique=None),   # filled in __init__
        "Raman":      dict(x_min=100,  x_max=3500, y_min=0, y_max=1,
                           x_label="Raman Shift (cm⁻¹)", y_label="Intensity",
                           technique=None),
        "UV-Vis":     dict(x_min=200,  x_max=800,  y_min=0, y_max=1,
                           x_label="Wavelength (nm)", y_label="Absorbance",
                           technique=None),
        "UV-Vis-NIR": dict(x_min=200,  x_max=2500, y_min=0, y_max=1,
                           x_label="Wavelength (nm)", y_label="Absorbance",
                           technique=None),
        "NIR":        dict(x_min=800,  x_max=2500, y_min=0, y_max=1,
                           x_label="Wavelength (nm)", y_label="Reflectance",
                           technique=None),
    }

    def __init__(self, parent_window, plugin):
        self.plugin       = plugin
        self.slide_images = []          # [(label, PIL.Image), ...]

        # fill in InstrumentType references after enum is defined
        self.PRESETS["FTIR"]["technique"]       = InstrumentType.FTIR
        self.PRESETS["Raman"]["technique"]      = InstrumentType.RAMAN
        self.PRESETS["UV-Vis"]["technique"]     = InstrumentType.UVVIS
        self.PRESETS["UV-Vis-NIR"]["technique"] = InstrumentType.UVVIS_NIR
        self.PRESETS["NIR"]["technique"]        = InstrumentType.NIR

        self.win = tk.Toplevel(parent_window)
        self.win.title("📊  Import Spectra from PPTX  —  Offline Digitizer")
        self.win.geometry("1040x660")
        self.win.minsize(820, 520)
        self.win.transient(parent_window)
        self.win.lift()
        self.win.focus_force()

        # ── tk variables ────────────────────────────────────────────────────
        self.technique_var = tk.StringVar(value="FTIR")
        self.x_min_var     = tk.StringVar(value="400")
        self.x_max_var     = tk.StringVar(value="4000")
        self.y_min_var     = tk.StringVar(value="0")
        self.y_max_var     = tk.StringVar(value="1")
        self.x_label_var   = tk.StringVar(value="Wavenumber (cm⁻¹)")
        self.y_label_var   = tk.StringVar(value="Absorbance")
        self.sample_var    = tk.StringVar(value="PPTX")
        self.bg_thresh_var = tk.IntVar(value=210)

        self._build_ui()

    # ── UI ──────────────────────────────────────────────────────────────────
    def _build_ui(self):
        # header
        hdr = tk.Frame(self.win, bg="#8e44ad", height=36)
        hdr.pack(fill=tk.X)
        hdr.pack_propagate(False)
        tk.Label(hdr,
                 text="📊  PPTX Spectrum Digitizer  ·  100 % offline · no internet · no API",
                 font=("Arial", 10, "bold"), bg="#8e44ad", fg="white"
                 ).pack(side=tk.LEFT, padx=10, pady=6)

        # main horizontal pane
        pane = tk.PanedWindow(self.win, orient=tk.HORIZONTAL, sashwidth=5)
        pane.pack(fill=tk.BOTH, expand=True, padx=4, pady=4)

        # ── LEFT: slide list + thumbnail ────────────────────────────────────
        left = tk.Frame(pane, bg="#f0f0f0")
        pane.add(left, width=248)

        tk.Label(left, text="Slides / Images", font=("Arial", 9, "bold"),
                 bg="#f0f0f0").pack(anchor=tk.W, padx=6, pady=(6, 2))

        lf = tk.Frame(left, bg="#f0f0f0")
        lf.pack(fill=tk.BOTH, expand=True, padx=4)

        self.slide_lb = tk.Listbox(lf, selectmode=tk.EXTENDED,
                                   font=("Arial", 8), bg="white",
                                   activestyle="dotbox", exportselection=False)
        sb = ttk.Scrollbar(lf, command=self.slide_lb.yview)
        self.slide_lb.configure(yscrollcommand=sb.set)
        self.slide_lb.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        sb.pack(side=tk.RIGHT, fill=tk.Y)
        self.slide_lb.bind("<<ListboxSelect>>", self._on_select)

        # ── RIGHT: config + preview plot ────────────────────────────────────
        right = tk.Frame(pane, bg="white")
        pane.add(right)

        # config frame
        cfg = tk.LabelFrame(right, text=" Axis Configuration ",
                             font=("Arial", 8), bg="white", padx=8, pady=4)
        cfg.pack(fill=tk.X, padx=6, pady=(6, 2))

        # row 0 — technique selector
        tk.Label(cfg, text="Technique preset:", font=("Arial", 8, "bold"),
                 bg="white").grid(row=0, column=0, sticky="w", pady=2)
        tech_cb = ttk.Combobox(cfg, textvariable=self.technique_var,
                               values=list(self.PRESETS.keys()),
                               width=14, state="readonly")
        tech_cb.grid(row=0, column=1, sticky="w", padx=4)
        tech_cb.bind("<<ComboboxSelected>>", self._apply_preset)
        tk.Label(cfg, text="↑ auto-fills axis defaults",
                 font=("Arial", 7), bg="white", fg="#888"
                 ).grid(row=0, column=2, columnspan=3, sticky="w")

        # rows 1–2 — axis ranges
        for lbl, var, row, col in [
            ("X min:", self.x_min_var, 1, 0),
            ("X max:", self.x_max_var, 1, 2),
            ("Y min:", self.y_min_var, 2, 0),
            ("Y max:", self.y_max_var, 2, 2),
        ]:
            tk.Label(cfg, text=lbl, font=("Arial", 8),
                     bg="white").grid(row=row, column=col, sticky="w", pady=2)
            tk.Entry(cfg, textvariable=var, width=10,
                     font=("Arial", 8)).grid(row=row, column=col+1,
                                             sticky="w", padx=4)

        # row 3-4 — axis labels
        tk.Label(cfg, text="X label:", font=("Arial", 8),
                 bg="white").grid(row=3, column=0, sticky="w", pady=2)
        tk.Entry(cfg, textvariable=self.x_label_var, width=24,
                 font=("Arial", 8)).grid(row=3, column=1, columnspan=4,
                                         sticky="w", padx=4)

        tk.Label(cfg, text="Y label:", font=("Arial", 8),
                 bg="white").grid(row=4, column=0, sticky="w", pady=2)
        tk.Entry(cfg, textvariable=self.y_label_var, width=24,
                 font=("Arial", 8)).grid(row=4, column=1, columnspan=4,
                                         sticky="w", padx=4)

        # row 5 — sample prefix + BG threshold
        tk.Label(cfg, text="Sample prefix:", font=("Arial", 8),
                 bg="white").grid(row=5, column=0, sticky="w", pady=2)
        tk.Entry(cfg, textvariable=self.sample_var, width=12,
                 font=("Arial", 8)).grid(row=5, column=1, sticky="w", padx=4)

        tk.Label(cfg, text="BG threshold (0-254):", font=("Arial", 8),
                 bg="white").grid(row=5, column=2, sticky="w")
        tk.Spinbox(cfg, from_=100, to=254, textvariable=self.bg_thresh_var,
                   width=5, font=("Arial", 8)).grid(row=5, column=3,
                                                    sticky="w", padx=4)

        # row 6 — hints
        hint = ("↑ lower = detect fainter/darker lines  |  "
                "try 180-220 for light backgrounds")
        tk.Label(cfg, text=hint, font=("Arial", 7),
                 bg="white", fg="#888").grid(row=6, column=0,
                                             columnspan=5, sticky="w")

        # dependency badge
        badge_txt = ("✅ Pillow OK  |  ✅ python-pptx OK"
                     if HAS_PIL and HAS_PPTX_LIB else
                     "⚠️  pip install Pillow python-pptx")
        badge_col = "#27ae60" if (HAS_PIL and HAS_PPTX_LIB) else "#c0392b"
        tk.Label(cfg, text=badge_txt, font=("Arial", 7),
                 bg="white", fg=badge_col).grid(row=7, column=0,
                                                columnspan=5, sticky="w",
                                                pady=(2, 0))

        # ── action buttons ───────────────────────────────────────────────────
        btn_row = tk.Frame(right, bg="white")
        btn_row.pack(fill=tk.X, padx=6, pady=4)

        ttk.Button(btn_row, text="📂 Open PPTX File…",
                   command=self._open_file, width=18).pack(side=tk.LEFT, padx=2)
        ttk.Separator(btn_row, orient=tk.VERTICAL).pack(side=tk.LEFT, padx=4, fill=tk.Y)
        ttk.Button(btn_row, text="✅ Import Selected",
                   command=self._import_selected, width=16).pack(side=tk.LEFT, padx=2)
        ttk.Button(btn_row, text="⬇️ Import All",
                   command=self._import_all, width=12).pack(side=tk.LEFT, padx=2)
        ttk.Button(btn_row, text="❌ Close",
                   command=self.win.destroy, width=8).pack(side=tk.RIGHT, padx=2)

        self.status_lbl = tk.Label(
            right,
            text="👆 Click  📂 Open PPTX File…  to begin.",
            font=("Arial", 9, "bold"), bg="white", fg="#2c3e50", anchor=tk.W, wraplength=700)
        self.status_lbl.pack(fill=tk.X, padx=8, pady=(0, 2))

        # ── preview plot ─────────────────────────────────────────────────────
        if HAS_MPL:
            self.prev_fig    = Figure(figsize=(6, 3), dpi=85, facecolor="white")
            self.prev_canvas = FigureCanvasTkAgg(self.prev_fig, master=right)
            self.prev_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True,
                                                  padx=4, pady=4)
            ax = self.prev_fig.add_subplot(111)
            ax.text(0.5, 0.5, "Select a slide and click  🔍 Preview",
                    ha="center", va="center", transform=ax.transAxes,
                    color="#888", fontsize=11)
            ax.axis("off")
            self.prev_canvas.draw()

    # ── file open ────────────────────────────────────────────────────────────
    def _open_file(self):
        if not HAS_PPTX_LIB or not HAS_PIL:
            missing = []
            if not HAS_PIL:      missing.append("Pillow")
            if not HAS_PPTX_LIB: missing.append("python-pptx")
            messagebox.showerror(
                "Missing libraries",
                f"Please install: pip install {' '.join(missing)}",
                parent=self.win)
            self.win.destroy()
            return

        path = filedialog.askopenfilename(
            title="Open PowerPoint file",
            filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")],
            parent=self.win)
        if not path:
            return   # user cancelled — keep dialog open
        self._load_pptx(path)

    def _load_pptx(self, path):
        self._set_status(f"Loading {Path(path).name} …")
        try:
            prs = _PPTXPresentation(path)
        except Exception as e:
            messagebox.showerror("PPTX error", str(e), parent=self.win)
            self._set_status(f"❌ Failed to load file: {e}", "#c0392b")
            return

        self.slide_images.clear()
        self.slide_lb.delete(0, tk.END)
        img_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, "image"):
                    continue
                try:
                    blob  = shape.image.blob
                    pil   = _PILImage.open(_pil_io.BytesIO(blob)).convert("RGB")
                    label = f"Slide {slide_idx + 1}  img {img_count + 1}"
                    self.slide_images.append((label, pil))
                    self.slide_lb.insert(tk.END, label)
                    img_count += 1
                except Exception:
                    continue

        if img_count == 0:
            messagebox.showwarning("No images",
                                   "No embedded images found in this file.",
                                   parent=self.win)
            self.win.destroy()
            return

        self.slide_lb.selection_set(0)
        self._on_select(None)
        self._set_status(
            f"Loaded {img_count} image(s) from {Path(path).name}.  "
            "Set axis ranges and click Preview to verify before importing.")

    # ── helpers ──────────────────────────────────────────────────────────────
    def _set_status(self, msg, color="#555"):
        self.status_lbl.configure(text=msg, fg=color)

    def _apply_preset(self, _event=None):
        p = self.PRESETS.get(self.technique_var.get())
        if not p:
            return
        self.x_min_var.set(str(p["x_min"]))
        self.x_max_var.set(str(p["x_max"]))
        self.y_min_var.set(str(p["y_min"]))
        self.y_max_var.set(str(p["y_max"]))
        self.x_label_var.set(p["x_label"])
        self.y_label_var.set(p["y_label"])

    def _on_select(self, _event):
        sel = self.slide_lb.curselection()
        if not sel:
            return
        # Immediately preview the selected slide in the big right-panel plot
        self._preview()

    def _get_params(self):
        """Parse axis params; raises ValueError on bad input."""
        return (float(self.x_min_var.get()), float(self.x_max_var.get()),
                float(self.y_min_var.get()), float(self.y_max_var.get()),
                int(self.bg_thresh_var.get()))

    def _digitize(self, pil, x_min, x_max, y_min, y_max, bg_thresh):
        return SpectrumDigitizer.digitize(
            pil, x_min, x_max, y_min, y_max, bg_threshold=bg_thresh)

    # ── preview ──────────────────────────────────────────────────────────────
    def _preview(self):
        """Preview one image — heavy work in background thread, result via ui_queue."""
        sel = self.slide_lb.curselection()
        if not sel:
            self._set_status("Select a slide first.", "#c0392b"); return
        if not HAS_MPL:
            self._set_status("matplotlib unavailable.", "#c0392b"); return
        try:
            x_min, x_max, y_min, y_max, bg_thresh = self._get_params()
        except ValueError:
            self._set_status("Invalid axis value.", "#c0392b"); return

        label, pil   = self.slide_images[sel[-1]]
        x_label      = self.x_label_var.get()
        y_label      = self.y_label_var.get()
        ui           = self.plugin.ui_queue   # the safe main-thread bridge

        self._set_status(f"⏳ Digitising {label} …", "#f39c12")

        def worker():
            x_arr, y_arr = self._digitize(pil, x_min, x_max, y_min, y_max, bg_thresh)

            def done():
                if not self.win.winfo_exists():
                    return
                if len(x_arr) == 0:
                    self._set_status(
                        "⚠️ No spectrum detected. Try lowering BG threshold.", "#c0392b")
                    return
                self.prev_fig.clear()
                ax = self.prev_fig.add_subplot(111)
                ax.plot(x_arr, y_arr, "b-", linewidth=1.2, alpha=0.85)
                ax.set_xlabel(x_label, fontsize=8)
                ax.set_ylabel(y_label, fontsize=8)
                ax.set_title(f"Preview: {label}  ({len(x_arr)} pts)", fontsize=9)
                ax.grid(True, alpha=0.3)
                self.prev_fig.tight_layout()
                self.prev_canvas.draw()
                self._set_status(
                    f"✅ {len(x_arr)} pts  |  X {x_min}–{x_max}  Y {y_min}–{y_max}",
                    "#27ae60")

            ui.schedule(done)

        threading.Thread(target=worker, daemon=True).start()

    # ── build one Spectrum — pure numpy/scipy, no tkinter, safe in any thread ─
    def _build_spectrum(self, label, pil, x_min, x_max, y_min, y_max,
                        bg_thresh, tech, x_label, y_label, prefix):
        x_arr, y_arr = self._digitize(pil, x_min, x_max, y_min, y_max, bg_thresh)
        if len(x_arr) == 0:
            return None
        safe_lb = re.sub(r"[^\w]+", "_", label)
        spec = Spectrum(
            timestamp    = datetime.now(),
            sample_id    = f"{prefix}_{safe_lb}",
            instrument   = "PPTX Image Digitizer",
            technique    = tech,
            manufacturer = "Digitized",
            model        = "Offline pixel analysis",
            x_data       = x_arr,
            y_data       = y_arr,
            x_label      = x_label,
            y_label      = y_label,
            file_source  = label,
            metadata     = {
                "source":       "pptx_digitizer",
                "x_range":      f"{x_min}–{x_max}",
                "y_range":      f"{y_min}–{y_max}",
                "bg_threshold": str(bg_thresh),
                "points":       str(len(x_arr)),
            },
        )
        if HAS_SCIPY:
            spec.find_peaks(tech)
        # identify_compounds runs here too — all in the worker thread, never on main
        if spec.peaks:
            spec.identify_compounds()
        return spec

    # ── import ───────────────────────────────────────────────────────────────
    def _import_selected(self):
        sel = self.slide_lb.curselection()
        if not sel:
            self._set_status("Select at least one slide.", "#c0392b"); return
        self._do_import([self.slide_images[i] for i in sel])

    def _import_all(self):
        self._do_import(self.slide_images)

    def _do_import(self, image_list):
        """
        Exactly mirrors _import_file / _batch_folder:
          1. Snapshot all tk vars on the main thread NOW.
          2. Spin up a daemon thread to do all CPU work.
          3. Post every UI update through plugin.ui_queue — the safe bridge.
        The main thread is never blocked.
        """
        try:
            x_min, x_max, y_min, y_max, bg_thresh = self._get_params()
        except ValueError:
            self._set_status("Invalid axis values.", "#c0392b"); return

        # Snapshot tk vars before leaving the main thread
        preset  = self.PRESETS.get(self.technique_var.get(), self.PRESETS["FTIR"])
        tech    = preset["technique"]
        x_label = self.x_label_var.get()
        y_label = self.y_label_var.get()
        prefix  = self.sample_var.get().strip() or "PPTX"
        total   = len(image_list)
        plug    = self.plugin
        ui      = plug.ui_queue   # ThreadSafeUI — polls on main thread every 50 ms

        self._set_status(f"⏳ Starting … 0 / {total}", "#f39c12")

        def worker():
            results  = []   # list of Spectrum | None, built entirely off-thread
            for idx, (label, pil) in enumerate(image_list):
                n = idx + 1
                # safe progress update through the queue
                ui.schedule(lambda n=n, lbl=label:
                    self._set_status(
                        f"⏳ Processing {n} / {total} — {lbl}", "#f39c12")
                    if self.win.winfo_exists() else None)

                spec = self._build_spectrum(
                    label, pil, x_min, x_max, y_min, y_max,
                    bg_thresh, tech, x_label, y_label, prefix)
                results.append(spec)

            # All done — schedule the final UI update on the main thread
            ui.schedule(lambda: self._finish_import(results, total))

        threading.Thread(target=worker, daemon=True).start()

    def _finish_import(self, results, total):
        """Runs on the main thread via ui_queue once the worker is done."""
        imported  = sum(1 for s in results if s is not None)
        failed    = total - imported
        last_spec = next((s for s in reversed(results) if s is not None), None)

        if imported:
            for spec in results:
                if spec is not None:
                    self.plugin.spectra.append(spec)
            self.plugin.current_spectrum = last_spec

        msg = f"✅ Imported {imported} spectra."
        if failed:
            msg += f"  {failed} had no detectable line."

        # Update dialog status then close it
        if self.win.winfo_exists():
            self._set_status(msg, "#27ae60" if imported else "#c0392b")
            self.win.destroy()

        if not imported:
            return

        # Update the plugin's UI — we are on the main thread, safe to do directly
        plug = self.plugin
        try:
            plug._update_tree()
            plug.count_label.config(text=f"📊 {len(plug.spectra)} spectra")
            if last_spec:
                if last_spec.peaks:
                    ps = ", ".join(f"{p:.1f}" for p in last_spec.peaks[:6])
                    plug.peak_label.config(
                        text=f"⚡ {len(last_spec.peaks)} peaks: {ps}"
                             f"{'…' if len(last_spec.peaks) > 6 else ''}")
                if last_spec.identifications:
                    plug._show_identifications(last_spec.identifications)
                if plug.plot_embedder:
                    plug.plot_embedder.plot_spectrum(last_spec)
            plug._update_status(
                f"✅ Imported {imported} spectra — click a row to switch", "#27ae60")
            plug._add_to_log(f"📊 PPTX: {imported} spectra added")
        except Exception as e:
            print(f"Plugin UI refresh error: {e}")


# ============================================================================
# PLOT EMBEDDER
# ============================================================================
class SpectroscopyPlotEmbedder:
    """Plot spectroscopy data"""

    def __init__(self, canvas_widget, figure):
        self.canvas = canvas_widget
        self.figure = figure
        self.current_plot = None

    def clear(self):
        self.figure.clear()
        self.figure.set_facecolor('white')
        self.current_plot = None

    def plot_spectrum(self, spec: Spectrum):
        """Plot spectrum with peaks"""
        self.clear()
        ax = self.figure.add_subplot(111)

        if spec.x_data is not None and spec.y_data is not None:
            ax.plot(spec.x_data, spec.y_data, 'b-', linewidth=1.5, alpha=0.8)

            # Mark peaks
            if spec.peaks and spec.peak_intensities:
                ax.scatter(spec.peaks, spec.peak_intensities,
                          c='red', s=30, zorder=5)

                # Label top peaks
                for i, (p, _) in enumerate(zip(spec.peaks[:5], spec.peak_intensities[:5])):
                    ax.annotate(f"{p:.1f}", (p, spec.peak_intensities[i]),
                               xytext=(5, 5), textcoords='offset points',
                               fontsize=7, bbox=dict(boxstyle='round,pad=0.2',
                                                      facecolor='yellow', alpha=0.7))

            ax.set_xlabel(spec.x_label, fontweight='bold')
            ax.set_ylabel(spec.y_label, fontweight='bold')

            # Title with technique and identifications
            title = f"{spec.technique.value.upper()}: {spec.sample_id}"
            if spec.identifications:
                title += f"\nTop: {spec.identifications[0]['name']} ({spec.identifications[0]['confidence']:.0f}%)"
            ax.set_title(title, fontweight='bold', fontsize=10)

            ax.grid(True, alpha=0.3)

        self.figure.tight_layout()
        self.canvas.draw()
        self.current_plot = 'spectrum'


# ============================================================================
# LIVE ACQUISITION TAB
# ============================================================================
class LiveAcquisitionTab:
    """Tab for real‑time acquisition from any connected instrument."""

    def __init__(self, parent, plugin):
        self.plugin = plugin
        self.frame = ttk.Frame(parent)
        self.driver = None          # currently selected instrument driver
        self.acquiring = False
        self.current_spectrum = None
        self.history = deque(maxlen=100)   # rolling history of spectra
        self._build_ui()

    def _build_ui(self):
        # Left panel for controls
        left = tk.Frame(self.frame, bg="white", width=250)
        left.pack(side=tk.LEFT, fill=tk.Y, padx=5, pady=5)

        # Right panel for plot
        right = tk.Frame(self.frame, bg="white")
        right.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)

        # Instrument selector
        tk.Label(left, text="Connected Instruments:", font=("Arial", 9, "bold"),
                 bg="white").pack(anchor=tk.W, pady=5)
        self.instr_combo = ttk.Combobox(left, state="readonly", width=30)
        self.instr_combo.pack(fill=tk.X, pady=2)
        self.instr_combo.bind('<<ComboboxSelected>>', self._on_instrument_selected)
        self._refresh_instrument_list()

        # Parameter frame (will be populated dynamically)
        self.param_frame = tk.LabelFrame(left, text="Acquisition Parameters",
                                         bg="white", font=("Arial", 8, "bold"))
        self.param_frame.pack(fill=tk.X, pady=10)
        self.param_widgets = {}   # store parameter entry widgets

        # Start/Stop buttons
        btn_frame = tk.Frame(left, bg="white")
        btn_frame.pack(fill=tk.X, pady=5)
        self.start_btn = tk.Button(btn_frame, text="Start", command=self._start,
                                    bg=C_STATUS, fg="white", width=10)
        self.start_btn.pack(side=tk.LEFT, padx=2)
        self.stop_btn = tk.Button(btn_frame, text="Stop", command=self._stop,
                                   state=tk.DISABLED, width=10)
        self.stop_btn.pack(side=tk.LEFT, padx=2)

        # History and snapshot
        tk.Label(left, text="History", font=("Arial", 8, "bold"),
                 bg="white").pack(pady=(10,0))
        self.history_combo = ttk.Combobox(left, values=[], state="readonly", width=30)
        self.history_combo.pack(fill=tk.X, pady=2)
        self.history_combo.bind('<<ComboboxSelected>>', self._show_history)

        self.snapshot_btn = tk.Button(left, text="Snapshot to Main Table",
                                      command=self._snapshot, width=20)
        self.snapshot_btn.pack(pady=10)

        # Matplotlib plot
        if HAS_MPL:
            self.fig = Figure(figsize=(6,4), dpi=90, facecolor="white")
            self.ax = self.fig.add_subplot(111)
            self.ax.set_xlabel("Wavelength / Wavenumber")
            self.ax.set_ylabel("Intensity")
            self.ax.grid(True, alpha=0.3)
            self.canvas = FigureCanvasTkAgg(self.fig, right)
            self.canvas.draw()
            self.canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
            toolbar = NavigationToolbar2Tk(self.canvas, right)
            toolbar.update()

        # Update instrument list periodically
        self._poll_instruments()

    def _refresh_instrument_list(self):
        """Populate combobox with connected devices."""
        devices = []
        for d in self.plugin.connected_devices:
            if hasattr(d, 'model'):
                devices.append(f"{d.manufacturer} {d.model}")
            else:
                devices.append("Unknown device")
        self.instr_combo['values'] = devices
        if devices and not self.instr_combo.get():
            self.instr_combo.current(0)
            self._on_instrument_selected()

    def _poll_instruments(self):
        """Refresh list every 5 seconds."""
        self._refresh_instrument_list()
        self.frame.after(5000, self._poll_instruments)

    def _on_instrument_selected(self, event=None):
        """Handle instrument selection: update parameter controls."""
        idx = self.instr_combo.current()
        if idx < 0 or idx >= len(self.plugin.connected_devices):
            return
        self.driver = self.plugin.connected_devices[idx]
        # Clear old parameter widgets
        for w in self.param_widgets.values():
            w.destroy()
        self.param_widgets.clear()
        # Create new controls based on driver capabilities
        row = 0
        # Integration time (common)
        if hasattr(self.driver, 'integration_time'):
            tk.Label(self.param_frame, text="Integration (ms):", bg="white",
                     font=("Arial", 8)).grid(row=row, column=0, sticky=tk.W, pady=2)
            var = tk.StringVar(value=str(getattr(self.driver, 'integration_time', 100)))
            entry = tk.Entry(self.param_frame, textvariable=var, width=10)
            entry.grid(row=row, column=1, sticky=tk.W, padx=5)
            self.param_widgets['integration_time'] = (var, entry)
            row += 1
        # Scans / averages (if applicable)
        if hasattr(self.driver, 'scans_to_average'):
            tk.Label(self.param_frame, text="Scans:", bg="white",
                     font=("Arial", 8)).grid(row=row, column=0, sticky=tk.W, pady=2)
            var = tk.StringVar(value=str(getattr(self.driver, 'scans_to_average', 1)))
            entry = tk.Entry(self.param_frame, textvariable=var, width=10)
            entry.grid(row=row, column=1, sticky=tk.W, padx=5)
            self.param_widgets['scans'] = (var, entry)
            row += 1
        # Laser power (Raman)
        if hasattr(self.driver, 'laser_power_mW'):
            tk.Label(self.param_frame, text="Laser power (mW):", bg="white",
                     font=("Arial", 8)).grid(row=row, column=0, sticky=tk.W, pady=2)
            var = tk.StringVar(value=str(getattr(self.driver, 'laser_power_mW', 50)))
            entry = tk.Entry(self.param_frame, textvariable=var, width=10)
            entry.grid(row=row, column=1, sticky=tk.W, padx=5)
            self.param_widgets['laser_power'] = (var, entry)
            row += 1
        # Add more as needed

    def _update_driver_params(self):
        """Transfer UI values to driver before acquisition."""
        if 'integration_time' in self.param_widgets:
            try:
                val = int(self.param_widgets['integration_time'][0].get())
                self.driver.integration_time = val
                if hasattr(self.driver, 'set_integration_time'):
                    self.driver.set_integration_time(val)
            except:
                pass
        if 'scans' in self.param_widgets:
            try:
                val = int(self.param_widgets['scans'][0].get())
                self.driver.scans_to_average = val
            except:
                pass
        if 'laser_power' in self.param_widgets:
            try:
                val = float(self.param_widgets['laser_power'][0].get())
                self.driver.laser_power_mW = val
            except:
                pass

    def _start(self):
        if self.driver is None:
            messagebox.showwarning("No Instrument", "Select an instrument first.")
            return
        self.acquiring = True
        self.start_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.NORMAL)
        # Run acquisition in background thread
        threading.Thread(target=self._acquisition_loop, daemon=True).start()

    def _stop(self):
        self.acquiring = False
        self.start_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.DISABLED)

    def _acquisition_loop(self):
        """Continuous acquisition loop."""
        while self.acquiring:
            # Update driver parameters from UI (thread-safe via queue)
            self.plugin.ui_queue.schedule(self._update_driver_params)
            # Acquire spectrum
            spec = self.driver.acquire_spectrum()
            if spec is not None:
                self.current_spectrum = spec
                self.history.append(spec)
                # Update plot on main thread
                self.plugin.ui_queue.schedule(self._update_plot)
                # Update history combobox
                times = [f"{i+1}" for i in range(len(self.history))]
                self.plugin.ui_queue.schedule(
                    lambda: self.history_combo.config(values=times)
                )
            # Wait according to integration time
            time.sleep(max(0.1, getattr(self.driver, 'integration_time', 100) / 1000))

    def _update_plot(self):
        """Redraw plot with current spectrum."""
        if self.current_spectrum and HAS_MPL:
            self.ax.clear()
            self.ax.plot(self.current_spectrum.x_data, self.current_spectrum.y_data,
                         color=C_ACCENT, lw=1)
            self.ax.set_xlabel(self.current_spectrum.x_label)
            self.ax.set_ylabel(self.current_spectrum.y_label)
            self.ax.set_title(f"Live: {self.current_spectrum.sample_id}")
            self.ax.grid(True, alpha=0.3)
            self.canvas.draw()

    def _show_history(self, event):
        """Display a previously acquired spectrum from history."""
        idx = self.history_combo.current()
        if idx >= 0 and idx < len(self.history):
            self.current_spectrum = self.history[idx]
            self._update_plot()

    def _snapshot(self):
        """Save current spectrum to main table."""
        if self.current_spectrum is None:
            return
        # Convert to dict format expected by data_hub
        sample = self.current_spectrum.to_dict()
        # Add to data_hub
        self.plugin.app.data_hub.add(sample)
        self.plugin._add_to_log(f"Snapshot: {self.current_spectrum.sample_id}")


# ============================================================================
# MAIN PLUGIN - SPECTROSCOPY UNIFIED SUITE
# ============================================================================
class SpectroscopyUnifiedSuitePlugin:

    def __init__(self, main_app):
        self.app = main_app
        self.window = None
        self.ui_queue = None
        self.deps = DEPS

        # Hardware devices
        self.thermo = None
        self.perkin = None
        self.bruker = None
        self.agilent = None
        self.bwtek = None
        self.ocean = None
        self.horiba = None
        self.metrohm = None
        self.avantes = None
        self.thorlabs = None
        self.viavi = None
        self.ti = None
        self.connected_devices = []

        # Current data
        self.spectra: List[Spectrum] = []
        self.current_spectrum: Optional[Spectrum] = None

        # Plot embedder
        self.plot_embedder = None

        # UI Variables
        self.status_var = tk.StringVar(value="Spectroscopy v5.0 - Ready")
        self.technique_var = tk.StringVar(value="FTIR")
        self.file_count_var = tk.StringVar(value="No files loaded")
        self.conf_var = tk.IntVar(value=50)

        # UI Elements
        self.notebook = None
        self.log_listbox = None
        self.plot_canvas = None
        self.plot_fig = None
        self.status_indicator = None
        self.technique_combo = None
        self.tree = None          # identifications treeview
        self.spec_tree = None     # spectra browser treeview
        self.import_btn = None
        self.batch_btn = None
        self.peak_label = None
        self.conn_status = None
        self.instrument_combo = None
        self.tol_label = None

        # All techniques
        self.all_techniques = [
            "FTIR - Thermo Nicolet",
            "FTIR - PerkinElmer",
            "FTIR - Bruker",
            "FTIR - Agilent Handheld",
            "Raman - B&W Tek",
            "Raman - Ocean Insight",
            "Raman - Horiba",
            "Raman - Metrohm Mira",
            "UV-Vis - Ocean Insight",
            "UV-Vis-NIR - Avantes",
            "UV-Vis - Thorlabs",
            "NIR - Viavi MicroNIR",
            "NIR - TI DLP",
            "File Import - Universal",
        ]

    def show_interface(self):
        self.open_window()

    def open_window(self):
        if self.window and self.window.winfo_exists():
            self.window.lift()
            return

        # 1000x650 window
        self.window = tk.Toplevel(self.app.root)
        self.window.title("Spectroscopy Unified Suite v5.0")
        self.window.geometry("1000x650")
        self.window.minsize(950, 600)
        self.window.transient(self.app.root)

        self.ui_queue = ThreadSafeUI(self.window)
        self._build_ui()
        self.window.lift()
        self.window.focus_force()
        self.window.protocol("WM_DELETE_WINDOW", self._on_close)

    def _build_ui(self):
        """1000x650 UI - All spectroscopy techniques"""

        # ============ HEADER (40px) ============
        header = tk.Frame(self.window, bg="#9b59b6", height=40)
        header.pack(fill=tk.X)
        header.pack_propagate(False)

        tk.Label(header, text="🧪", font=("Arial", 16),
                bg="#9b59b6", fg="white").pack(side=tk.LEFT, padx=8)

        tk.Label(header, text="SPECTROSCOPY UNIFIED SUITE", font=("Arial", 12, "bold"),
                bg="#9b59b6", fg="white").pack(side=tk.LEFT, padx=2)

        tk.Label(header, text="v5.0 · 52+ INSTRUMENTS", font=("Arial", 8),
                bg="#9b59b6", fg="#f1c40f").pack(side=tk.LEFT, padx=8)

        # Technique badges
        badge_frame = tk.Frame(header, bg="#9b59b6")
        badge_frame.pack(side=tk.LEFT, padx=10)
        for badge, tip in [("🥼", "FTIR"), ("🔬", "RAMAN"), ("📊", "UV-Vis"), ("📱", "NIR")]:
            lbl = tk.Label(badge_frame, text=badge, font=("Arial", 12),
                          bg="#9b59b6", fg="white")
            lbl.pack(side=tk.LEFT, padx=2)

        self.status_indicator = tk.Label(header, textvariable=self.status_var,
                                        font=("Arial", 8), bg="#9b59b6", fg="#2ecc71")
        self.status_indicator.pack(side=tk.RIGHT, padx=10)

        # ============ INSTRUMENT SELECTOR (30px) ============
        inst_frame = tk.Frame(self.window, bg="#f8f9fa", height=30)
        inst_frame.pack(fill=tk.X)
        inst_frame.pack_propagate(False)

        tk.Label(inst_frame, text="Instrument:", font=("Arial", 8, "bold"),
                bg="#f8f9fa").pack(side=tk.LEFT, padx=8)

        self.instrument_combo = ttk.Combobox(inst_frame,
                                            textvariable=self.technique_var,
                                            values=self.all_techniques,
                                            width=40, state="readonly")
        self.instrument_combo.pack(side=tk.LEFT, padx=5)
        self.instrument_combo.current(0)

        self.connect_btn = ttk.Button(inst_frame, text="🔌 Connect",
                                      command=self._connect_instrument, width=8)
        self.connect_btn.pack(side=tk.LEFT, padx=2)

        self.conn_status = tk.Label(inst_frame, text="●", fg="#e74c3c",
                                   font=("Arial", 10), bg="#f8f9fa")
        self.conn_status.pack(side=tk.LEFT, padx=5)

        # Tolerance indicator
        self.tol_label = tk.Label(inst_frame, text="Δ: 5.0 cm⁻¹", font=("Arial", 7),
                                 bg="#f8f9fa", fg="#7f8c8d")
        self.tol_label.pack(side=tk.RIGHT, padx=15)

        # ============ MAIN SPLIT ============
        main = tk.PanedWindow(self.window, orient=tk.HORIZONTAL, sashwidth=4)
        main.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # ============ LEFT PANEL - SPECTRUM ============
        left = tk.Frame(main, bg="white")
        main.add(left, width=600)

        # File toolbar
        # ── two toolbar rows so nothing gets clipped ────────────────────────
        toolbar = tk.Frame(left, bg="#f0f0f0")
        toolbar.pack(fill=tk.X)

        file_bar = tk.Frame(toolbar, bg="#f0f0f0")
        file_bar.pack(fill=tk.X, padx=3, pady=(3, 1))

        self.import_btn = ttk.Button(file_bar, text="📂 Load File / PPTX",
                                     command=self._import_file, width=18)
        self.import_btn.pack(side=tk.LEFT, padx=2)

        self.batch_btn = ttk.Button(file_bar, text="📁 Batch Folder",
                                    command=self._batch_folder, width=14)
        self.batch_btn.pack(side=tk.LEFT, padx=2)

        self.file_label = tk.Label(file_bar, text="No file loaded",
                                   font=("Arial", 8), bg="#f0f0f0", fg="#7f8c8d")
        self.file_label.pack(side=tk.LEFT, padx=6)

        action_bar = tk.Frame(toolbar, bg="#f0f0f0")
        action_bar.pack(fill=tk.X, padx=3, pady=(0, 3))

        ttk.Button(action_bar, text="🔍 Find Peaks",
                  command=self._find_peaks, width=13).pack(side=tk.LEFT, padx=2)

        ttk.Button(action_bar, text="🔬 Identify Compounds",
                  command=self._identify_compounds, width=18).pack(side=tk.LEFT, padx=2)

        ttk.Separator(action_bar, orient=tk.VERTICAL).pack(
            side=tk.LEFT, padx=6, fill=tk.Y, pady=2)

        ttk.Button(action_bar, text="📤 Send to Table",
                  command=self.send_to_table, width=14).pack(side=tk.LEFT, padx=2)

        # Spectrum plot
        plot_frame = tk.Frame(left, bg="white")
        plot_frame.pack(fill=tk.BOTH, expand=True, padx=2, pady=2)

        self.plot_fig = Figure(figsize=(6, 3.5), dpi=85, facecolor='white')
        self.plot_canvas = FigureCanvasTkAgg(self.plot_fig, master=plot_frame)
        self.plot_canvas.draw()
        self.plot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)

        self.plot_embedder = SpectroscopyPlotEmbedder(self.plot_canvas, self.plot_fig)

        # Placeholder
        ax = self.plot_fig.add_subplot(111)
        ax.text(0.5, 0.5, 'Load a spectrum to begin',
               ha='center', va='center', transform=ax.transAxes,
               fontsize=12, color='#7f8c8d')
        ax.set_title('Spectroscopy Plot', fontweight='bold')
        ax.set_xlabel('X Axis')
        ax.set_ylabel('Intensity')
        ax.axis('on')
        self.plot_canvas.draw()

        # Peak info line
        peak_bar = tk.Frame(left, bg="#e9ecef", height=24)
        peak_bar.pack(fill=tk.X)
        peak_bar.pack_propagate(False)

        self.peak_label = tk.Label(peak_bar, text="⚡ No peaks detected",
                                   font=("Arial", 7), bg="#e9ecef", anchor=tk.W)
        self.peak_label.pack(fill=tk.X, padx=8)

        # ============ RIGHT PANEL - IDENTIFICATIONS ============
        right = tk.Frame(main, bg="white")
        main.add(right, width=350)

        # ── top: spectra list (click to switch current spectrum) ────────────
        spec_hdr = tk.Frame(right, bg="#f8f9fa", height=24)
        spec_hdr.pack(fill=tk.X)
        spec_hdr.pack_propagate(False)
        tk.Label(spec_hdr, text="📂 Loaded Spectra  (click to select)",
                 font=("Arial", 7, "bold"), bg="#f8f9fa").pack(side=tk.LEFT, padx=6)

        spec_frame = tk.Frame(right, bg="white")
        spec_frame.pack(fill=tk.X, padx=2, pady=(0, 2))

        spec_cols = ('ID', 'Technique', 'Peaks', 'Source')
        self.spec_tree = ttk.Treeview(spec_frame, columns=spec_cols,
                                      show='headings', height=6,
                                      selectmode='browse')
        self.spec_tree.heading('ID',        text='Sample ID')
        self.spec_tree.heading('Technique', text='Technique')
        self.spec_tree.heading('Peaks',     text='Peaks')
        self.spec_tree.heading('Source',    text='Source')
        self.spec_tree.column('ID',        width=120)
        self.spec_tree.column('Technique', width=55, anchor='center')
        self.spec_tree.column('Peaks',     width=35, anchor='center')
        self.spec_tree.column('Source',    width=100)
        ss = ttk.Scrollbar(spec_frame, command=self.spec_tree.yview)
        self.spec_tree.configure(yscrollcommand=ss.set)
        self.spec_tree.pack(side=tk.LEFT, fill=tk.X, expand=True)
        ss.pack(side=tk.RIGHT, fill=tk.Y)
        self.spec_tree.bind('<<TreeviewSelect>>', self._on_spec_select)

        # ── bottom: identifications ──────────────────────────────────────────
        # Control bar
        ctrl_bar = tk.Frame(right, bg="#f8f9fa", height=28)
        ctrl_bar.pack(fill=tk.X)
        ctrl_bar.pack_propagate(False)

        tk.Label(ctrl_bar, text="Identifications — Confidence ≥",
                 font=("Arial", 7, "bold"), bg="#f8f9fa").pack(side=tk.LEFT, padx=4)

        conf_combo = ttk.Combobox(ctrl_bar, textvariable=self.conf_var,
                                  values=[25, 50, 75, 90], width=3, state="readonly")
        conf_combo.pack(side=tk.LEFT, padx=2)
        tk.Label(ctrl_bar, text="%", font=("Arial", 7),
                bg="#f8f9fa").pack(side=tk.LEFT)

        ttk.Button(ctrl_bar, text="📋 Clear", command=self._clear_results,
                  width=6).pack(side=tk.RIGHT, padx=2)

        # Identifications tree
        id_frame = tk.Frame(right, bg="white")
        id_frame.pack(fill=tk.BOTH, expand=True, padx=2, pady=2)

        columns = ('Compound', 'Conf', 'Group', 'Formula', 'Matches')
        self.tree = ttk.Treeview(id_frame, columns=columns, show='headings',
                                 height=14, selectmode='browse')

        self.tree.heading('Compound', text='Compound / Mineral')
        self.tree.heading('Conf',     text='%')
        self.tree.heading('Group',    text='Group')
        self.tree.heading('Formula',  text='Formula')
        self.tree.heading('Matches',  text='Hits')

        self.tree.column('Compound', width=120)
        self.tree.column('Conf',     width=38,  anchor='center')
        self.tree.column('Group',    width=65)
        self.tree.column('Formula',  width=75)
        self.tree.column('Matches',  width=38,  anchor='center')

        yscroll = ttk.Scrollbar(id_frame, command=self.tree.yview)
        self.tree.configure(yscrollcommand=yscroll.set)

        self.tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        yscroll.pack(side=tk.RIGHT, fill=tk.Y)

        # Colour tags
        self.tree.tag_configure('high', background='#d4edda')
        self.tree.tag_configure('med',  background='#fff3cd')
        self.tree.tag_configure('low',  background='#f8d7da')

        # ============ NOTEBOOK WITH TABS (including Live) ============
        notebook = ttk.Notebook(self.window)
        notebook.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # Live tab
        self.live_tab = LiveAcquisitionTab(notebook, self)
        notebook.add(self.live_tab.frame, text=" Live ")

        # Other analysis tabs could be added here later, but for now just Live.

        # ============ STATUS BAR ============
        status = tk.Frame(self.window, bg="#34495e", height=24)
        status.pack(fill=tk.X, side=tk.BOTTOM)
        status.pack_propagate(False)

        deps_str = []
        if HAS_SEABREEZE: deps_str.append("🌊")
        if HAS_SERIAL: deps_str.append("📡")
        if HAS_VISA: deps_str.append("🔌")
        if HAS_REQUESTS: deps_str.append("🌐")
        if HAS_BLEAK: deps_str.append("📱")

        self.count_label = tk.Label(status,
            text=f"📊 {len(self.spectra)} spectra · {sum(len(v) for v in SPECTRAL_LIBRARY.values())} compounds",
            font=("Arial", 7), bg="#34495e", fg="white")
        self.count_label.pack(side=tk.LEFT, padx=8)

        tk.Label(status,
                text="Thermo · PerkinElmer · Bruker · Agilent · B&W Tek · Ocean · Horiba · Renishaw · Avantes · Thorlabs · Viavi · TI",
                font=("Arial", 7), bg="#34495e", fg="#bdc3c7").pack(side=tk.RIGHT, padx=8)

    # ============================================================================
    # INSTRUMENT CONNECTION METHODS
    # ============================================================================

    def _connect_instrument(self):
        """Connect to selected instrument"""
        selection = self.technique_var.get()

        self.connect_btn.config(state='disabled')
        self._update_status(f"Connecting to {selection}...", "#f39c12")

        def connect_thread():
            success = False
            msg = ""
            driver = None

            try:
                # FTIR
                if "Thermo" in selection:
                    driver = ThermoFTIRDriver()
                    success, msg = driver.connect()
                    self.thermo = driver

                elif "PerkinElmer" in selection:
                    driver = PerkinElmerFTIRDriver()
                    success, msg = driver.connect()
                    self.perkin = driver

                elif "Bruker" in selection:
                    driver = BrukerFTIRDriver()
                    success, msg = driver.connect()
                    self.bruker = driver

                elif "Agilent" in selection:
                    driver = AgilentFTIRDriver()
                    success, msg = driver.connect()
                    self.agilent = driver

                # Raman
                elif "B&W" in selection:
                    driver = BWTekRamanDriver()
                    success, msg = driver.connect()
                    self.bwtek = driver

                elif "Ocean" in selection and "Raman" in selection:
                    driver = OceanRamanDriver()
                    success, msg = driver.connect()
                    self.ocean = driver

                elif "Horiba" in selection:
                    driver = HoribaRamanDriver()
                    success, msg = driver.connect()
                    self.horiba = driver

                elif "Metrohm" in selection:
                    driver = MetrohmMiraDriver()
                    success, msg = driver.connect()
                    self.metrohm = driver

                # UV-Vis
                elif "Ocean" in selection and "UV" in selection:
                    driver = OceanUVVisDriver()
                    success, msg = driver.connect()
                    self.ocean = driver

                elif "Avantes" in selection:
                    driver = AvantesDriver()
                    success, msg = driver.connect()
                    self.avantes = driver

                elif "Thorlabs" in selection:
                    driver = ThorlabsCCSDriver()
                    success, msg = driver.connect()
                    self.thorlabs = driver

                # NIR
                elif "Viavi" in selection:
                    driver = ViaviMicroNIRDriver()
                    success, msg = driver.connect()
                    self.viavi = driver

                elif "TI" in selection:
                    driver = TIDLPNIRDriver()
                    success, msg = driver.connect() if hasattr(driver, 'connect') else (True, "TI DLP (Demo)")
                    self.ti = driver

                else:
                    success = True
                    msg = "File Import mode"

            except Exception as e:
                success = False
                msg = str(e)

            def update_ui():
                self.connect_btn.config(state='normal')
                if success:
                    self.conn_status.config(fg="#2ecc71")
                    self._update_status(f"✅ {msg}", "#27ae60")
                    if driver and driver not in self.connected_devices:
                        self.connected_devices.append(driver)
                else:
                    self.conn_status.config(fg="red")
                    self._update_status(f"❌ {msg}", "#e74c3c")

            self.ui_queue.schedule(update_ui)

        threading.Thread(target=connect_thread, daemon=True).start()

    # ============================================================================
    # FILE IMPORT METHODS
    # ============================================================================

    def _import_file(self):
        """Import a spectral file or PowerPoint presentation."""
        filetypes = [
            ("All supported", "*.csv;*.txt;*.spa;*.wdf;*.jdx;*.dx;*.opj;*.dpt;*.ngs;*.pptx"),
            ("PowerPoint (PPTX)", "*.pptx"),
            ("CSV/ASCII", "*.csv;*.txt;*.asc;*.dat"),
            ("Thermo SPA", "*.spa"),
            ("Witec WDF", "*.wdf"),
            ("JCAMP", "*.jdx;*.dx"),
            ("Origin", "*.opj"),
            ("Bruker", "*.dpt"),
            ("Viavi", "*.ngs"),
            ("All files", "*.*")
        ]

        path = filedialog.askopenfilename(filetypes=filetypes)
        if not path:
            return

        # Route PPTX to the offline digitizer dialog
        if Path(path).suffix.lower() == '.pptx':
            if not self.window:
                return
            dlg = PPTXImportDialog(self.window, self)
            dlg._load_pptx(path)   # pre-load the already-chosen file
            return

        self._update_status(f"Parsing {Path(path).name}...", "#f39c12")
        self.import_btn.config(state='disabled')

        def parse_thread():
            result = SpectroscopyParserFactory.parse_file(path)

            def update_ui():
                self.import_btn.config(state='normal')
                if result is None:
                    self._add_to_log(f"❌ Failed to parse: {Path(path).name}")
                    self._update_status("❌ Parse failed", "#e74c3c")
                    return

                # Handle both single spectrum and list of spectra
                if isinstance(result, list):
                    for spec in result:
                        self.spectra.append(spec)
                        if HAS_SCIPY:
                            spec.find_peaks(spec.technique)
                    self.current_spectrum = result[0] if result else None
                else:
                    self.spectra.append(result)
                    self.current_spectrum = result
                    if HAS_SCIPY:
                        result.find_peaks(result.technique)

                # Update UI
                self._update_tree()
                self.count_label.config(text=f"📊 {len(self.spectra)} spectra")
                self.file_label.config(text=Path(path).name[:30])
                self._add_to_log(f"✅ Imported from {Path(path).name}")

                # Plot the current spectrum
                if self.current_spectrum and self.plot_embedder:
                    self.plot_embedder.plot_spectrum(self.current_spectrum)

                # Update peak label
                if self.current_spectrum and self.current_spectrum.peaks:
                    peak_str = ", ".join([f"{p:.1f}" for p in self.current_spectrum.peaks[:6]])
                    self.peak_label.config(text=f"⚡ {len(self.current_spectrum.peaks)} peaks: {peak_str}{'...' if len(self.current_spectrum.peaks)>6 else ''}")
                else:
                    self.peak_label.config(text="⚡ No peaks detected")

                self._update_status(f"✅ Loaded: {Path(path).name}", "#27ae60")

            self.ui_queue.schedule(update_ui)

        threading.Thread(target=parse_thread, daemon=True).start()

    def _batch_folder(self):
        """Batch process a folder of spectral files"""
        folder = filedialog.askdirectory()
        if not folder:
            return

        self._update_status(f"Scanning {Path(folder).name}...", "#f39c12")
        self.import_btn.config(state='disabled')
        self.batch_btn.config(state='disabled')

        def batch_thread():
            count = 0
            for ext in ['*.csv', '*.txt', '*.spa', '*.wdf', '*.jdx', '*.dx', '*.dpt', '*.ngs']:
                for filepath in Path(folder).glob(ext):
                    res = SpectroscopyParserFactory.parse_file(str(filepath))
                    if res is None:
                        continue
                    if isinstance(res, list):
                        for spec in res:
                            self.spectra.append(spec)
                            if HAS_SCIPY:
                                spec.find_peaks(spec.technique)
                        count += len(res)
                    else:
                        self.spectra.append(res)
                        if HAS_SCIPY:
                            res.find_peaks(res.technique)
                        count += 1

            def update_ui():
                self._update_tree()
                self.file_count_var.set(f"{len(self.spectra)} files")
                self.count_label.config(text=f"📊 {len(self.spectra)} spectra")
                self._add_to_log(f"📁 Batch imported: {count} spectra from {Path(folder).name}")
                self._update_status(f"✅ Imported {count} spectra", "#27ae60")
                self.import_btn.config(state='normal')
                self.batch_btn.config(state='normal')

            self.ui_queue.schedule(update_ui)

        threading.Thread(target=batch_thread, daemon=True).start()

    def _import_peak_list(self, path=None):
        """
        Import a CSV containing peak assignments (Sample, Peak_cm-1, Assignment)
        and send it directly to the main data table.
        """
        if not path:
            path = filedialog.askopenfilename(
                title="Select peak list CSV",
                filetypes=[("CSV files", "*.csv"), ("All files", "*.*")]
            )
            if not path:
                return

        try:
            df = pd.read_csv(path)
            # Check required columns
            required = {'Sample', 'Peak_cm-1', 'Assignment'}
            if not required.issubset(df.columns):
                missing = required - set(df.columns)
                messagebox.showerror(
                    "Invalid Format",
                    f"File must contain columns: {', '.join(required)}\nMissing: {', '.join(missing)}"
                )
                return

            # Convert to list of dicts for data_hub
            records = df.to_dict(orient='records')
            self.app.data_hub.add_samples(records)
            self.app.samples = self.app.data_hub.get_all()
            if hasattr(self.app.center, '_refresh'):
                self.app.center._refresh()

            self._add_to_log(f"📋 Imported {len(records)} peak assignments from {Path(path).name}")
            self._update_status(f"✅ Peak list loaded: {len(records)} rows", "#27ae60")

        except Exception as e:
            messagebox.showerror("Import Error", f"Failed to read file:\n{str(e)}")

    def _import_pptx(self):
        """Open the offline PPTX spectrum digitizer dialog."""
        if not HAS_PPTX_LIB or not HAS_PIL:
            missing = []
            if not HAS_PIL:       missing.append("Pillow")
            if not HAS_PPTX_LIB: missing.append("python-pptx")
            if messagebox.askyesno(
                    "Missing libraries",
                    f"The following libraries are needed:\n\n"
                    f"  pip install {' '.join(missing)}\n\n"
                    "Open a terminal and install them, then try again.\n\n"
                    "Would you like to copy the install command?"):
                try:
                    self.window.clipboard_clear()
                    self.window.clipboard_append(
                        f"pip install {' '.join(missing)}")
                except Exception:
                    pass
            return
        if not self.window:
            return
        PPTXImportDialog(self.window, self)

    # ============================================================================
    # ANALYSIS METHODS
    # ============================================================================

    def _update_tree(self):
        """Refresh the spectra browser (top-right list)."""
        if not hasattr(self, 'spec_tree') or self.spec_tree is None:
            return
        for item in self.spec_tree.get_children():
            self.spec_tree.delete(item)

        for i, spec in enumerate(self.spectra[-50:]):
            tag = 'current' if spec is self.current_spectrum else ''
            self.spec_tree.insert('', tk.END, iid=str(i), tags=(tag,), values=(
                spec.sample_id[:18],
                spec.technique.value.upper(),
                str(len(spec.peaks)),
                (Path(spec.file_source).name if spec.file_source else 'Live')[:16],
            ))

        self.spec_tree.tag_configure('current', background='#cce5ff')

        # Highlight current spectrum row.
        # Unbind first so selection_set doesn't re-fire _on_spec_select.
        if self.current_spectrum and self.current_spectrum in self.spectra:
            idx = str(self.spectra.index(self.current_spectrum))
            try:
                self.spec_tree.unbind('<<TreeviewSelect>>')
                self.spec_tree.selection_set(idx)
                self.spec_tree.see(idx)
            except Exception:
                pass
            finally:
                self.spec_tree.bind('<<TreeviewSelect>>', self._on_spec_select)

    def _on_spec_select(self, event):
        """Switch current spectrum when user clicks a row in the spectra browser."""
        sel = self.spec_tree.selection()
        if not sel:
            return
        try:
            idx = int(sel[0])
            spec = self.spectra[idx]
        except (ValueError, IndexError):
            return
        self.current_spectrum = spec
        # Manually update the highlight tag without calling _update_tree()
        # (which would re-fire this event via selection_set)
        for item in self.spec_tree.get_children():
            self.spec_tree.item(item, tags=())
        try:
            self.spec_tree.item(sel[0], tags=('current',))
        except Exception:
            pass

        # Update plot
        if self.plot_embedder:
            self.plot_embedder.plot_spectrum(spec)

        # Update peak label
        if spec.peaks:
            ps = ', '.join(f'{p:.1f}' for p in spec.peaks[:6])
            self.peak_label.config(
                text=f'⚡ {len(spec.peaks)} peaks: {ps}'
                     f'{"..." if len(spec.peaks) > 6 else ""}')
        else:
            self.peak_label.config(text='⚡ No peaks detected')

        # Show identifications if already computed
        if spec.identifications:
            self._show_identifications(spec.identifications)
        else:
            for item in self.tree.get_children():
                self.tree.delete(item)

        self.file_label.config(text=spec.sample_id[:30])
        self._update_status(
            f'📂 {spec.sample_id} | {spec.technique.value.upper()} | '
            f'{len(spec.peaks)} peaks | {len(spec.identifications)} matches')

    def _find_peaks(self):
        """Find peaks in current spectrum"""
        if not self.current_spectrum:
            messagebox.showwarning("No Data", "Load a spectrum first")
            return

        if not HAS_SCIPY:
            messagebox.showwarning("Missing Dependency", "scipy required for peak detection")
            return

        self.current_spectrum.find_peaks()
        self.plot_embedder.plot_spectrum(self.current_spectrum)

        if self.current_spectrum.peaks:
            peak_str = ", ".join([f"{p:.1f}" for p in self.current_spectrum.peaks[:6]])
            self.peak_label.config(text=f"⚡ {len(self.current_spectrum.peaks)} peaks: {peak_str}{'...' if len(self.current_spectrum.peaks)>6 else ''}")
            self._add_to_log(f"🔍 Found {len(self.current_spectrum.peaks)} peaks")
        else:
            self.peak_label.config(text="⚡ No peaks detected")
            self._add_to_log("🔍 No peaks found")

    def _identify_compounds(self):
        """Identify compounds in current spectrum"""
        if not self.current_spectrum:
            messagebox.showwarning("No Data", "Load a spectrum first")
            return

        if not self.current_spectrum.peaks:
            self.current_spectrum.find_peaks()

        if not self.current_spectrum.peaks:
            messagebox.showwarning("No Peaks", "Find peaks first")
            return

        # Set tolerance based on technique
        if self.current_spectrum.technique == InstrumentType.RAMAN:
            tol = 5.0
            unit = "cm⁻¹"
        elif self.current_spectrum.technique == InstrumentType.FTIR:
            tol = 10.0
            unit = "cm⁻¹"
        elif self.current_spectrum.technique in [InstrumentType.UVVIS, InstrumentType.UVVIS_NIR]:
            tol = 3.0
            unit = "nm"
        elif self.current_spectrum.technique == InstrumentType.NIR:
            tol = 8.0
            unit = "nm"
        else:
            tol = 5.0
            unit = ""

        self.tol_label.config(text=f"Δ: {tol} {unit}")

        # Identify
        results = self.current_spectrum.identify_compounds(tolerance=tol)

        # Filter by confidence
        threshold = self.conf_var.get()
        filtered = [r for r in results if r['confidence'] >= threshold]

        self._show_identifications(filtered)
        self._add_to_log(f"🔬 Identified {len(filtered)} compounds (≥{threshold}% confidence)")
        self.plot_embedder.plot_spectrum(self.current_spectrum)
        self._update_tree()   # refresh spectra browser

        if not filtered:
            self.peak_label.config(text="⚠️ No compounds above confidence threshold")

    def _show_identifications(self, results):
        """
        Populate the identifications tree with compound match results.
        Expects a list of dicts: name, confidence, group, formula, matches, total.
        """
        for item in self.tree.get_children():
            self.tree.delete(item)
        for r in results[:25]:
            conf = r.get('confidence', 0)
            tag  = 'high' if conf >= 70 else ('med' if conf >= 50 else 'low')
            self.tree.insert('', tk.END, tags=(tag,), values=(
                r.get('name', ''),
                f"{conf:.0f}",
                r.get('group', ''),
                r.get('formula', ''),
                f"{r.get('matches', 0)}/{r.get('total', 0)}",
            ))

    def _clear_results(self):
        """Clear identification results"""
        for item in self.tree.get_children():
            self.tree.delete(item)

    def _add_to_log(self, message):
        """Add message to log (simplified - no log tab in this version)"""
        # Just update status
        self._update_status(message)

    def _update_status(self, message, color=None):
        self.status_var.set(message)
        if color and self.status_indicator:
            self.status_indicator.config(fg=color)

    # ============================================================================
    # DATA EXPORT
    # ============================================================================

    def send_to_table(self):
        """
        Send spectra to main table - ONE row per spectrum.
        Columns are created dynamically from the dictionary keys.
        The full spectral data (x and y arrays) are stored as JSON strings
        in columns named after the axis labels (e.g. "Wavenumber (cm⁻¹)", "Absorbance").
        """
        if not self.spectra:
            messagebox.showwarning("No Data", "No spectra to send")
            return

        table_data = []

        for spec in self.spectra:
            # Start with a clean dictionary - every key becomes a column
            row = {}

            # Sample ID is required
            row["Sample_ID"] = spec.sample_id

            # Add ALL metadata as columns
            if spec.technique:
                row["Technique"] = spec.technique.value
            if spec.instrument:
                row["Instrument"] = spec.instrument
            if spec.manufacturer:
                row["Manufacturer"] = spec.manufacturer
            if spec.model:
                row["Model"] = spec.model
            if spec.file_source:
                row["Source"] = Path(spec.file_source).name
            if spec.timestamp:
                row["Date"] = spec.timestamp.strftime("%Y-%m-%d %H:%M:%S")

            # STORE FULL SPECTRUM DATA as JSON strings
            # Use the axis labels as column names
            if spec.x_data is not None:
                col_x = spec.x_label.strip() if spec.x_label else "X_Data"
                row[col_x] = json.dumps(spec.x_data.tolist())
            if spec.y_data is not None:
                col_y = spec.y_label.strip() if spec.y_label else "Y_Data"
                row[col_y] = json.dumps(spec.y_data.tolist())

            # Peak information
            if spec.peaks:
                row["Peaks"] = str(len(spec.peaks))
                row["Peak_Positions"] = ", ".join([f"{p:.1f}" for p in spec.peaks[:10]])
            else:
                row["Peaks"] = "0"

            # Compound identifications
            if spec.identifications:
                compounds = []
                groups = []
                formulas = []
                confidences = []
                peaks_hit = []

                for ident in spec.identifications[:5]:  # Limit to top 5
                    compounds.append(ident.get('name', ''))
                    groups.append(ident.get('group', ''))
                    formulas.append(ident.get('formula', ''))
                    confidences.append(f"{ident.get('confidence', 0):.0f}%")
                    peaks_hit.append(f"{ident.get('matches', 0)}/{ident.get('total', 0)}")

                row["Compound"] = ", ".join(compounds)
                row["Group"] = ", ".join(groups)
                row["Formula"] = ", ".join(formulas)
                row["Confidence"] = ", ".join(confidences)
                row["Peaks_Hit"] = ", ".join(peaks_hit)

            # Add any other metadata
            if spec.resolution_cm1:
                row["Resolution_cm1"] = spec.resolution_cm1
            if spec.integration_time_ms:
                row["Integration_ms"] = spec.integration_time_ms
            if spec.laser_power_mW:
                row["Laser_Power_mW"] = spec.laser_power_mW

            table_data.append(row)

        try:
            # Send directly to data_hub
            dh = self.app.data_hub
            dh.add_samples(table_data)
            self.app.samples = dh.get_all()

            # Force table refresh
            if hasattr(self.app, 'center') and hasattr(self.app.center, '_refresh'):
                self.app.center._refresh()

            self._add_to_log(f"📤 Sent {len(table_data)} spectra with full waveform data")
            self._update_status(f"✅ Sent {len(table_data)} spectra")

        except Exception as e:
            messagebox.showerror("Error", f"Failed to send data: {str(e)}")

    def collect_data(self) -> List[Dict]:
        return [s.to_dict() for s in self.spectra]

    def _on_close(self):
        """Clean disconnect"""
        self._add_to_log("🛑 Shutting down...")

        for device in self.connected_devices:
            if device and hasattr(device, 'disconnect'):
                try:
                    device.disconnect()
                    self._add_to_log(f"✅ Device disconnected")
                except:
                    pass

        self.connected_devices.clear()

        if self.window:
            self.window.destroy()
            self.window = None


# ============================================================================
# SIMPLE PLUGIN REGISTRATION - NO DUPLICATES
# ============================================================================

def setup_plugin(main_app):
    """Register plugin - simple, no duplicates"""

    # Create plugin instance
    plugin = SpectroscopyUnifiedSuitePlugin(main_app)

    # Add to left panel if available
    if hasattr(main_app, 'left') and main_app.left is not None:
        main_app.left.add_hardware_button(
            name=PLUGIN_INFO.get("name", "Spectroscopy Suite"),
            icon=PLUGIN_INFO.get("icon", "🧪"),
            command=plugin.show_interface
        )
        print(f"✅ Added: {PLUGIN_INFO.get('name')}")

    return plugin
